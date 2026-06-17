from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pathlib import Path

src = Path('C:/Users/admin/Desktop/research/RL/manuscripts/competition_pitch_thesis_formula.pptx')
out = Path('C:/Users/admin/Desktop/research/RL/manuscripts/competition_pitch_thesis_formula_v3.pptx')
prs = Presentation(str(src))

DARK = RGBColor(44, 62, 80)
ACCENT = RGBColor(33, 118, 255)
LIGHT = RGBColor(245, 248, 252)
MID = RGBColor(99, 110, 123)
WHITE = RGBColor(255, 255, 255)


def apply_runs(shape, font_size=None, bold=None, color=None, align=None):
    if not hasattr(shape, 'text_frame'):
        return
    tf = shape.text_frame
    tf.word_wrap = True
    for p in tf.paragraphs:
        if align is not None:
            p.alignment = align
        for r in p.runs:
            if font_size is not None:
                r.font.size = Pt(font_size)
            if bold is not None:
                r.font.bold = bold
            if color is not None:
                r.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False,
                color=DARK, align=PP_ALIGN.LEFT, fill=None, line=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    else:
        box.fill.background()
    if line is not None:
        box.line.color.rgb = line
    else:
        box.line.fill.background()
    return box


def add_round_box(slide, left, top, width, height, text, font_size=16, fill=LIGHT, line=ACCENT, text_color=DARK):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    tf = shp.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.bold = True
    r.font.color.rgb = text_color
    return shp

# Slide 1 cover
s = prs.slides[0]
s.shapes[1].text = 'AI 自适应上肢康复机器人系统'
s.shapes[2].text = '面向基层康复的智能训练平台'
apply_runs(s.shapes[1], font_size=28, bold=True, color=DARK)
apply_runs(s.shapes[2], font_size=16, color=MID)
add_textbox(s, 1400000, 4300000, 9300000, 900000,
            '用强化学习驱动康复机器人，根据患者实时能力自动调节训练难度。',
            font_size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_textbox(s, 1800000, 5200000, 8500000, 700000,
            '目标场景：康复医院 / 基层医院 / 康复中心 / 社区与居家康复',
            font_size=14, color=MID, align=PP_ALIGN.CENTER)

# Slide 2 pain
s = prs.slides[1]
apply_runs(s.shapes[0], font_size=24, bold=True, color=DARK)
for idx in [1, 3, 5, 7]:
    apply_runs(s.shapes[idx], font_size=20, bold=True, color=ACCENT)
for idx in [2, 4, 6, 8]:
    apply_runs(s.shapes[idx], font_size=14, color=DARK)

# Slide 3 why now
s = prs.slides[2]
s.shapes[0].text = '为什么现在值得做？需求、政策、国产替代与 AI 落地窗口正在重叠'
s.shapes[2].text = '这不是“只有技术没有市场”的科研项目，而是康复需求增长、政策支持、国产替代与 AI 工程化同步出现的窗口期。'
s.shapes[13].text = '需求端\n脑卒中患者规模大、上肢功能障碍康复周期长，老龄化进一步推高长期训练需求。'
s.shapes[14].text = '政策端\n康复辅助器具创新持续获政策支持，康复服务纳入支付体系的趋势也在增强基层采购动力。'
s.shapes[15].text = '供给端\n强化学习、机器人控制与视觉感知正在从实验室走向工程化，使自适应康复训练第一次具备产品化基础。'
apply_runs(s.shapes[0], font_size=22, bold=True, color=DARK)
apply_runs(s.shapes[2], font_size=14, color=MID)
for idx in [13, 14, 15]:
    apply_runs(s.shapes[idx], font_size=13, color=DARK)

# Slide 4 customers
s = prs.slides[3]
s.shapes[0].text = '谁会先买单，我们先从哪里切入？'
s.shapes[12].text = '首批客户'
s.shapes[13].text = '二级及以上医院康复科、专业康复医院，采购决策相对明确，适合作为首批试点与品牌背书场景。'
s.shapes[15].text = '增长客户'
s.shapes[16].text = '康复中心、社区卫生服务中心和基层医疗机构更关注低成本、易部署与远程支持，是后续放量关键。'
s.shapes[18].text = '延展场景'
s.shapes[19].text = '养老机构、居家康复与区域康复网络可在后续通过平台化服务接入，形成长期服务收入。'
s.shapes[21].text = '切入路径'
s.shapes[22].text = '先做医院试点和示范案例，再复制到区域渠道和基层机构，最后拓展远程康复生态。'
apply_runs(s.shapes[0], font_size=22, bold=True, color=DARK)
for idx in [12, 15, 18, 21]:
    apply_runs(s.shapes[idx], font_size=18, bold=True, color=ACCENT)
for idx in [13, 16, 19, 22]:
    apply_runs(s.shapes[idx], font_size=13, color=DARK)

# Slide 5 what we do
s = prs.slides[4]
s.shapes[0].text = '01'
s.shapes[1].text = '一句话讲清楚我们做什么'
s.shapes[2].text = 'WHAT WE DO'
apply_runs(s.shapes[0], font_size=26, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
apply_runs(s.shapes[1], font_size=28, bold=True, color=DARK)
apply_runs(s.shapes[2], font_size=16, color=MID)
add_textbox(s, 1000000, 4300000, 10200000, 1000000,
            '我们用强化学习驱动康复机器人，根据患者实时能力自动调节训练难度。',
            font_size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_textbox(s, 1700000, 5400000, 8800000, 600000,
            '让训练从“固定程序”变成“按人动态适配”。',
            font_size=15, color=MID, align=PP_ALIGN.CENTER)

# Slide 6 core value
s = prs.slides[5]
s.shapes[0].text = '核心价值：不是替代设备，而是让康复训练更个性化、更标准化、更可复制'
s.shapes[1].shapes[0].text = '对患者'
s.shapes[1].shapes[2].text = '训练难度随能力变化自动调整，更容易把患者维持在“够得着、又有挑战”的有效训练区间，提升参与感和坚持度。'
s.shapes[2].shapes[0].text = '对治疗师'
s.shapes[2].shapes[2].text = '减少重复调参与主观判断负担，用量化数据辅助复盘和干预决策，让服务交付更标准化。'
s.shapes[3].shapes[0].text = '对机构'
s.shapes[3].shapes[2].text = '以更低门槛引入智能康复能力，同时兼顾示范展示、科研合作和后续远程康复延展空间。'
apply_runs(s.shapes[0], font_size=22, bold=True, color=DARK)
for gidx in [1, 2, 3]:
    apply_runs(s.shapes[gidx].shapes[0], font_size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    apply_runs(s.shapes[gidx].shapes[2], font_size=14, color=DARK, align=PP_ALIGN.CENTER)

# Slide 7 product form
s = prs.slides[6]
apply_runs(s.shapes[0], font_size=22, bold=True, color=DARK)
for idx in [10, 12, 14, 16]:
    apply_runs(s.shapes[idx], font_size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
for idx in [11, 13, 15, 17]:
    apply_runs(s.shapes[idx], font_size=13, color=DARK)

# Slide 8 traction
s = prs.slides[7]
s.shapes[0].text = '阶段成果：项目已跨过概念验证，正在向实机验证和产品化推进'
s.shapes[15].shapes[0].text = '当前进展'
s.shapes[15].shapes[1].text = '已完成仿真训练框架、关键算法设计和实机平台初步打通，具备继续产品化推进的基础。'
s.shapes[16].shapes[0].text = '算法验证'
s.shapes[16].shapes[1].text = '已完成多种模型结构对比与迭代，验证“自适应训练”并非概念，而是一条可持续优化的技术路线。'
s.shapes[17].shapes[0].text = '实机部署'
s.shapes[17].shapes[1].text = '已完成协作机器人、交互终端和视觉感知的系统联调，证明方案不只停留在纸面。'
s.shapes[18].shapes[0].text = '临床试点'
s.shapes[18].shapes[1].text = '下一步计划联合康复医院开展试点验证，补齐真实患者场景下的安全性与有效性证据。'
s.shapes[19].shapes[0].text = '产品化'
s.shapes[19].shapes[1].text = '后续将推进工程化、注册路径与渠道验证，把科研能力逐步转化为可销售产品。'
apply_runs(s.shapes[0], font_size=22, bold=True, color=DARK)
for idx in [15, 16, 17, 18, 19]:
    apply_runs(s.shapes[idx].shapes[0], font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    apply_runs(s.shapes[idx].shapes[1], font_size=12, color=DARK, align=PP_ALIGN.CENTER)

# Slide 9 flow
s = prs.slides[8]
s.shapes[0].text = '02'
s.shapes[1].text = '患者如何使用这套系统？'
s.shapes[2].text = 'TRAINING FLOW'
apply_runs(s.shapes[0], font_size=26, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
apply_runs(s.shapes[1], font_size=28, bold=True, color=DARK)
apply_runs(s.shapes[2], font_size=16, color=MID)
add_round_box(s, 1000000, 4300000, 2200000, 1000000, '患者开始训练', font_size=18)
add_round_box(s, 3600000, 4300000, 2200000, 1000000, '系统感知状态', font_size=18)
add_round_box(s, 6200000, 4300000, 2200000, 1000000, 'AI 动态调难度', font_size=18)
add_round_box(s, 8800000, 4300000, 2200000, 1000000, '输出评估报告', font_size=18)
add_textbox(s, 1400000, 5600000, 9000000, 500000,
            '把训练过程、训练难度和训练结果形成闭环。',
            font_size=14, color=MID, align=PP_ALIGN.CENTER)

# Slide 10 difficulty
s = prs.slides[9]
s.shapes[0].text = '为什么这套系统短期内不容易被复制？'
s.shapes[1].text = ''
s.shapes[4].text = '难点一'
s.shapes[5].text = '真正的壁垒不是“机器人能动”，而是系统能持续识别患者能力变化，并把难度控制在合适区间。'
s.shapes[7].text = '难点二'
s.shapes[8].text = '要让策略适应不同患者，就需要足够贴近真实能力差异的训练体系，而不是只靠少量实验演示。'
s.shapes[10].text = '难点三'
s.shapes[11].text = '从仿真到真实设备的落地，涉及感知、控制、安全和部署成本的协同优化，比单点算法创新更难。'
apply_runs(s.shapes[0], font_size=22, bold=True, color=DARK)
for idx in [4, 7, 10]:
    apply_runs(s.shapes[idx], font_size=17, bold=True, color=ACCENT)
for idx in [5, 8, 11]:
    apply_runs(s.shapes[idx], font_size=13, color=DARK)

# Slide 11 moat
s = prs.slides[10]
s.shapes[0].text = '三类核心能力，构成我们的长期护城河'
s.shapes[1].text = '真正难复制的不是一个术语，而是我们把“训练策略、训练体系、真实落地”三件事同时往前推进。'
s.shapes[14].text = '自适应控制'
s.shapes[15].text = '根据患者表现动态调难度，让系统逐步替代部分高频人工调参工作，形成产品核心价值。'
s.shapes[16].text = '虚拟患者训练'
s.shapes[17].text = '通过更丰富的仿真训练对象覆盖患者能力谱系，提升系统面对复杂康复场景时的稳定性和泛化能力。'
s.shapes[18].text = '实机部署能力'
s.shapes[19].text = '完成从仿真到真实设备的迁移与联调，更接近可交付产品，而不仅是实验室演示结果。'
s.shapes[20].text = '长期护城河'
apply_runs(s.shapes[0], font_size=22, bold=True, color=DARK)
apply_runs(s.shapes[1], font_size=13, color=MID, align=PP_ALIGN.CENTER)
for idx in [14, 16, 18, 20]:
    apply_runs(s.shapes[idx], font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
for idx in [15, 17, 19]:
    apply_runs(s.shapes[idx], font_size=12, color=DARK, align=PP_ALIGN.CENTER)

# Slide 12 competition
s = prs.slides[11]
s.shapes[0].text = '03'
s.shapes[1].text = '竞品对比：传统设备重执行，我们更强调自适应与可复制'
s.shapes[2].text = 'COMPETITION'
apply_runs(s.shapes[0], font_size=26, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
apply_runs(s.shapes[1], font_size=26, bold=True, color=DARK)
apply_runs(s.shapes[2], font_size=16, color=MID)
rows, cols = 6, 4
left, top, width, height = 700000, 3800000, 10700000, 2600000
shape = s.shapes.add_table(rows, cols, left, top, width, height)
tbl = shape.table
headers = ['维度', '进口康复机器人', '传统训练设备', '我们的方案']
for c, h in enumerate(headers):
    cell = tbl.cell(0, c)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENT
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.bold = True
            r.font.size = Pt(12)
            r.font.color.rgb = WHITE
body = [
    ['训练方式', '固定程序或阻抗控制为主', '预设模式为主', '可随患者表现动态调难度'],
    ['部署门槛', '价格高、部署复杂', '能力有限但较易部署', '追求低成本与可落地平衡'],
    ['数据反馈', '部分可量化', '反馈较弱', '训练过程与评估结果可沉淀'],
    ['目标客户', '大型医院为主', '基础训练场景', '医院试点后可下沉基层'],
    ['核心差异', '硬件强', '设备轻', 'AI 自适应 + 仿真训练 + 实机联调']
]
for r in range(1, rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = body[r - 1][c]
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(250, 252, 255) if r % 2 else RGBColor(242, 246, 251)
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for rr in p.runs:
                rr.font.size = Pt(10.5)
                rr.font.color.rgb = DARK

# Slide 13 business model
s = prs.slides[12]
s.shapes[0].text = '商业模式：硬件切入，软件订阅与服务收入形成长期价值'
g = s.shapes[1]
g.shapes[0].shapes[2].text = '硬件销售：整套康复机器人系统，定价约 15–30 万元/台，先以国产替代切入。'
g.shapes[1].shapes[2].text = '软件订阅：训练算法、数据分析与增值功能按年收费，形成持续复购。'
g.shapes[2].shapes[2].text = '培训服务：设备使用培训与康复方案设计培训，增强客户黏性。'
g.shapes[3].shapes[2].text = '远程康复与数据服务：打开长期服务收入和科研合作空间。'
apply_runs(s.shapes[0], font_size=22, bold=True, color=DARK)
for subgroup in [g.shapes[0], g.shapes[1], g.shapes[2], g.shapes[3]]:
    apply_runs(subgroup.shapes[1], font_size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    apply_runs(subgroup.shapes[2], font_size=12, color=DARK, align=PP_ALIGN.CENTER)

# Slide 14 GTM
s = prs.slides[13]
s.shapes[0].text = '市场进入路径：先试点验证，再区域复制，最后向基层与远程康复扩展'
g = s.shapes[1]
g.shapes[7].shapes[1].text = '1–2 年：联合 3–5 家医院做试点，形成真实案例、临床背书和产品打样。'
g.shapes[8].shapes[1].text = '2–3 年：围绕重点省份建立渠道合作，复制到二级以上医院和康复中心。'
g.shapes[9].shapes[1].text = '3–5 年：通过远程平台降低基层使用门槛，进入社区、基层机构与居家康复场景。'
g.shapes[10].shapes[1].text = '持续推进：沉淀训练数据与服务能力，逐步构建康复管理平台。'
g.shapes[11].shapes[1].text = '同步任务：推进注册、合作网络与品牌影响力，形成可持续增长飞轮。'
apply_runs(s.shapes[0], font_size=22, bold=True, color=DARK)
for idx in [7, 8, 9, 10, 11]:
    apply_runs(g.shapes[idx].shapes[0], font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    apply_runs(g.shapes[idx].shapes[1], font_size=12, color=DARK, align=PP_ALIGN.CENTER)

# Slide 15 team and ask
s = prs.slides[14]
s.shapes[0].text = '团队与合作诉求：我们更需要试点场景、产业资源与资金支持'
s.shapes[1].shapes[3].text = '当前基础'
s.shapes[2].shapes[12].text = '强化学习'
s.shapes[2].shapes[13].text = '机器人控制'
s.shapes[3].shapes[9].text = '计算机视觉'
s.shapes[3].shapes[10].text = '康复场景'
s.shapes[4].shapes[8].text = '渠道合作'
s.shapes[4].shapes[9].text = '医院试点'
s.shapes[5].text = '融资支持'
apply_runs(s.shapes[0], font_size=22, bold=True, color=DARK)
apply_runs(s.shapes[1].shapes[3], font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
for shp in [s.shapes[2].shapes[12], s.shapes[2].shapes[13], s.shapes[3].shapes[9], s.shapes[3].shapes[10], s.shapes[4].shapes[8], s.shapes[4].shapes[9], s.shapes[5]]:
    apply_runs(shp, font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_textbox(s, 1400000, 5200000, 4200000, 900000,
            '合作诉求：试点医院、康复中心、产业导师与渠道伙伴。',
            font_size=14, color=DARK, align=PP_ALIGN.CENTER)
add_textbox(s, 6500000, 5200000, 4200000, 900000,
            '融资诉求：用于工程化推进、临床试点、注册准备与市场验证。',
            font_size=14, color=DARK, align=PP_ALIGN.CENTER)

prs.save(str(out))
print(out)
