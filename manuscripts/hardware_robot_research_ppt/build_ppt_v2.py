from pathlib import Path
from urllib.request import urlretrieve
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

BASE = Path(__file__).resolve().parent
ASSETS = BASE / "assets_v2"
OUT = BASE / "康复桌面抓取目标微型移动机器人_运动机理分类调研_汇报版_v3_ESPROLL.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

NAVY = RGBColor(27, 48, 67)
TEAL = RGBColor(16, 139, 137)
TEXT = RGBColor(37, 47, 57)
MUTED = RGBColor(96, 108, 119)
LINE = RGBColor(215, 222, 228)
PALE = RGBColor(242, 247, 248)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(48, 124, 83)
ORANGE = RGBColor(208, 126, 38)
RED = RGBColor(180, 63, 70)
BLUE = RGBColor(43, 97, 150)
FONT = "Microsoft YaHei"

URL = {
    "toio": "https://toio.io/platform/cube/",
    "toio_spec": "https://toio.github.io/toio-spec/en/docs/hardware_other/",
    "toio_motor": "https://toio.github.io/toio-spec/en/docs/motor/",
    "epuck": "https://projects.gctronic.com/epuck2/e-puck2-flyer.pdf",
    "zooids": "https://shape.stanford.edu/research/swarm/SwarmUIs_no_copyright.pdf",
    "misaka": "https://arxiv.org/abs/2404.17125",
    "robot_soccer": "https://robot-soccer-kit.github.io/",
    "robot_soccer_photo": "https://hackaday.io/project/181818/gallery",
    "lgdx": "https://lgdxrobot.uk/lgdxrobot2/",
    "mecanum_accuracy": "https://www.jstage.jst.go.jp/article/jsmermd/2021/0/2021_1P1-L14/_article/-char/en",
    "esp_roll": "https://www.instructables.com/ESP-ROLL-Build-a-Spherical-Self-balancing-Robot-Wi/",
    "esp_roll_video": "https://www.youtube.com/watch?v=VuBfRYNQgw0",
    "esp_roll_files": "https://cults3d.com/en/3d-model/game/esp-roll-an-fpv-esp32-spherical-robot-rc-ball-bot-3d-design-stl-files",
    "esp_roll_img": "https://fbi.cults3d.com/uploaders/29670915/illustration-file/20d8febb-d2e8-47fa-bcd2-a203503484e2/ESP-ROLL-ESP32-FPV-Spherical-Robot-Ball-Bot.jpg",
    "sphere_paper": "https://research.aalto.fi/en/publications/mechanical-development-and-control-of-a-miniature-nonholonomic-sp/",
    "cellulo": "https://academia.skadge.org/publis/ozgur2017cellulo.pdf",
    "cellulo_photo": "https://actu.epfl.ch/news/a-robot-to-help-visually-impaired-schoolchildren-2/",
    "cellulo_control": "https://infoscience.epfl.ch/server/api/core/bitstreams/4dccb2b7-e91c-4fa9-9c95-c25ef41669c4/content",
    "modulo": "https://doi.org/10.1109/IROS47612.2022.9981983",
    "torus": "https://catalog.lib.kyushu-u.ac.jp/opac_download_md/7172240/7172240_AM.pdf",
}


def add_text(slide, text, x, y, w, h, size=18, color=TEXT, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, hyperlink=None,
             margin=0.0, font=FONT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    if hyperlink: r.hyperlink.address = hyperlink
    return box


def title(slide, text, subtitle=None):
    add_text(slide, text, 0.78, 0.58, 11.7, 0.52, 28, NAVY, True)
    if subtitle:
        add_text(slide, subtitle, 0.80, 1.12, 11.6, 0.32, 12.2, MUTED)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.78), Inches(1.53), Inches(11.78), Inches(0.025))
    line.fill.solid(); line.fill.fore_color.rgb = TEAL; line.line.fill.background()


def footer(slide, page, links=None):
    add_text(slide, f"运动机理分类调研  |  {page:02d}", 0.78, 7.15, 2.5, 0.18, 8.2, MUTED)
    if links:
        x = 7.15
        for label, url in links:
            w = max(0.72, min(1.75, len(label)*0.095+0.30))
            add_text(slide, label, x, 7.15, w, 0.18, 8.2, BLUE, hyperlink=url)
            x += w + 0.18


def ensure_asset(path, url):
    if not path.exists():
        urlretrieve(url, path)
    return path


def picture_contain(slide, path, x, y, w, h, link=None):
    path = str(path)
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w/iw, h/ih)
    pw, ph = iw*scale, ih*scale
    pic = slide.shapes.add_picture(path, Inches(x+(w-pw)/2), Inches(y+(h-ph)/2), Inches(pw), Inches(ph))
    if link: pic.click_action.hyperlink.address = link
    return pic


def specs(slide, rows, x=7.08, y=1.86, w=5.22, size=14.2, row_h=0.69):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(row_h*len(rows)+0.1))
    tf = box.text_frame
    tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (label, value) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        p.line_spacing = 1.0
        r1 = p.add_run(); r1.text = f"{label}  "
        r1.font.name = FONT; r1.font.size = Pt(size); r1.font.bold = True; r1.font.color.rgb = TEAL
        r2 = p.add_run(); r2.text = value
        r2.font.name = FONT; r2.font.size = Pt(size); r2.font.bold = False; r2.font.color.rgb = TEXT
    return box


# 1. Cover
s = prs.slides.add_slide(blank)
s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
add_text(s, "微型桌面小车", 0.92, 1.28, 6.2, 0.52, 18, TEAL, True)
add_text(s, "运动机理分类调研", 0.90, 1.93, 8.2, 0.82, 40, NAVY, True)
add_text(s, "面向康复桌面抓取目标的结构筛选", 0.92, 2.98, 6.4, 0.48, 20, TEXT)
add_text(s, "重点比较差速、全向轮、麦克纳姆、球形滚动与全向球驱动", 0.94, 3.72, 7.0, 0.62, 15, MUTED)
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.92), Inches(4.72), Inches(3.3), Inches(0.04))
line.fill.solid(); line.fill.fore_color.rgb = TEAL; line.line.fill.background()
add_text(s, "石慧峰  ·  2026.07", 0.94, 5.04, 3.8, 0.30, 13, MUTED)
picture_contain(s, ASSETS/"cellulo_body.jpg", 8.35, 1.28, 4.05, 4.75, URL["cellulo_photo"])
add_text(s, "代表性全向球驱动科研平台：Cellulo", 8.45, 6.14, 3.85, 0.25, 10, MUTED, align=PP_ALIGN.CENTER)

# 2. Requirements
s = prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
title(s, "场景要求决定底盘机理", "目标不是普通移动机器人，而是能够复现机械臂二维末端效果的小型抓取目标")
requirements = [
    ("尺寸", "优先30–80 mm；质量和外形应允许安全触碰或抓取"),
    ("运动学", "最好能够保持朝向横移，即独立控制 vx、vy 与角速度 ω"),
    ("平滑性", "5–20 mm/s低速不爬行；快速换向没有明显停顿、振动或轮滑"),
    ("控制", "外部20 Hz RL指令可持续接入，底层具有更高频速度/位置闭环"),
    ("精度", "不能把编码器分辨率当运动精度；目标是毫米级动态轨迹误差"),
]
y=1.92
for i,(a,b) in enumerate(requirements,1):
    add_text(s,f"{i:02d}",0.92,y,0.52,0.34,14,TEAL,True)
    add_text(s,a,1.64,y,1.20,0.34,17,NAVY,True)
    add_text(s,b,3.00,y,8.65,0.38,15,TEXT)
    y += 0.94
add_text(s,"筛选逻辑：先看机理能否全向，再看现有平台是否足够小、足够平滑、具有精度证据。",0.94,6.65,11.3,0.34,15,NAVY,True)
footer(s,2)

# 3. Taxonomy
s = prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
title(s, "五类主要运动机理", "从非完整约束到真正三自由度全向，结构自由度决定轨迹形态")
headers=["机理","典型控制量","保持朝向横移","低速主要问题","代表平台"]
widths=[2.15,2.20,1.78,3.22,2.48]
x0=0.78; y0=1.82; rh=0.75
x=x0
for h,w in zip(headers,widths):
    add_text(s,h,x,y0,w,0.34,12,NAVY,True,PP_ALIGN.CENTER)
    x+=w
rows=[
    ("双轮差速","v, ω","不能","转向停顿；非完整约束","toio / e-puck2"),
    ("三轮全向轮","vx, vy, ω","能","滚子切换与侧向打滑","Misaka / RSK"),
    ("四轮麦克纳姆","vx, vy, ω","能","滚子振动；横移里程计误差","LGDX / ROSbot"),
    ("球形滚动","方向 + 速度","通常不能严格实现","球壳惯性、低速漂移、换向滞后","Sphero / 学术球形原型"),
    ("全向球/特殊驱动","vx, vy, ω","能","制造复杂；平台多为科研原型","Cellulo / torus wheel"),
]
for i,row in enumerate(rows):
    y=y0+0.47+i*0.88
    ln=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x0),Inches(y-0.13),Inches(sum(widths)),Inches(0.012))
    ln.fill.solid(); ln.fill.fore_color.rgb=LINE; ln.line.fill.background()
    x=x0
    for j,(val,w) in enumerate(zip(row,widths)):
        c = GREEN if (j==2 and val=="能") else RED if (j==2 and val=="不能") else TEXT
        add_text(s,val,x+0.05,y,w-0.10,0.52,12.4,c,j==0,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
        x+=w
add_text(s,"只有后三类在运动学上具备机械臂式平面全向能力；其中全向球的连续接触最有利于低速平滑。",0.92,6.58,11.4,0.36,15,NAVY,True,PP_ALIGN.CENTER)
footer(s,3)

# 4. Differential
s = prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
title(s,"双轮差速驱动", "最容易做到微型化，但横向运动必须通过转向或圆弧完成")
picture_contain(s, ASSETS/"toio_body.png",0.82,1.86,5.55,4.55,URL["toio"])
add_text(s,"代表实物：Sony toio Core Cube（商用）",0.92,6.48,5.35,0.25,10,MUTED,align=PP_ALIGN.CENTER)
specs(s,[
    ("机理", "左右两轮独立调速，直接控制线速度v和角速度ω；属于非完整约束。"),
    ("代表", "商用toio 32×32×20 mm、30 g；科研e-puck2直径70 mm、130 g。"),
    ("速度", "toio最高350 mm/s；e-puck2最高约154 mm/s。"),
    ("控制", "BLE/Python目标点控制或ROS/外部视觉均成熟。"),
    ("精度", "toio无官方毫米实测；e-puck2的0.13 mm仅为电机步进分辨率。"),
    ("场景判断", "尺寸和成熟度最好，但不能保持朝向横移，无法复现机械臂二维末端轨迹。"),
])
add_text(s,"筛选：适合快速原型，不适合作为“机械臂效果”等效方案。",7.08,6.34,5.10,0.58,14,ORANGE,True)
footer(s,4,[("toio",URL["toio"]),("toio规格",URL["toio_spec"]),("e-puck2",URL["epuck"])])

# 5. Three omni
s = prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
title(s,"三轮全向轮驱动", "三个120°布置的全向轮能够直接产生任意平移与旋转")
picture_contain(s, ASSETS/"robot_soccer_omni.png",0.82,1.83,5.55,4.62,URL["robot_soccer_photo"])
add_text(s,"代表实物：Robot Soccer Kit（三全向轮科研/开源平台）",0.92,6.48,5.35,0.25,10,MUTED,align=PP_ALIGN.CENTER)
specs(s,[
    ("机理", "三个全向轮通常以120°布置，独立控制vx、vy、ω，是真正平面全向。"),
    ("代表", "Misaka直径100 mm、高50 mm；RSK直径约176 mm、710 g。"),
    ("速度", "两类平台公开最高平移速度均约200 mm/s。"),
    ("控制", "Arduino/XBee/纸面微点阵，或Python/蓝牙/顶部ArUco视觉。"),
    ("精度", "Misaka未报告整机轨迹误差；RSK视觉约3 mm/像素，未报告路径RMSE。"),
    ("场景判断", "运动学符合，但现有完整平台仍偏大，滚子接触会带来低速波动与侧滑。"),
])
add_text(s,"筛选：可作为第二候选机理，但需要进一步微型化并验证低速轨迹。",7.08,6.34,5.10,0.58,14,ORANGE,True)
footer(s,5,[("Misaka",URL["misaka"]),("Robot Soccer Kit",URL["robot_soccer"]),("照片",URL["robot_soccer_photo"])])

# 6. Mecanum
s = prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
title(s,"四轮麦克纳姆驱动", "四个带45°滚子的车轮可全向移动，但紧凑性和低速丝滑性较差")
picture_contain(s, ASSETS/"lgdx_mecanum.jpg",0.82,1.82,5.55,4.65,URL["lgdx"])
add_text(s,"代表实物：LGDXRobot2（开源ROS 2麦克纳姆平台）",0.92,6.48,5.35,0.25,10,MUTED,align=PP_ALIGN.CENTER)
specs(s,[
    ("机理", "四轮滚子轴通常为±45°，通过轮速组合独立生成vx、vy、ω。"),
    ("代表", "LGDXRobot2约240×240×134 mm；较小的MentorPi仍约212×171×147 mm、1.2 kg。"),
    ("速度", "同类ROS底盘常见0.6–1.0 m/s，远高于康复目标需求。"),
    ("控制", "STM32电机PID、编码器、IMU、ROS 2 cmd_vel与外部定位。"),
    ("精度", "小型产品通常无轨迹RMSE；横移与斜移时编码器里程计受轮滑影响明显。"),
    ("场景判断", "现有平台尺寸和质量过大；外露滚子、振动和侧向滑移不利于手部交互。"),
])
add_text(s,"筛选：运动学全向，但现有尺寸和机械接触特性不适合抓取目标。",7.08,6.34,5.10,0.58,14,RED,True)
footer(s,6,[("LGDXRobot2",URL["lgdx"]),("轮滑研究",URL["mecanum_accuracy"])])

# 7. Spherical
s = prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
title(s,"球形滚动机器人", "开源ESP-ROLL便于复现，但球体运动仍不等于机械臂式三自由度全向")
picture_contain(s, ensure_asset(ASSETS/"esp_roll_body.jpg", URL["esp_roll_img"]),0.82,1.84,5.55,4.60,URL["esp_roll_files"])
add_text(s,"代表实物：ESP-ROLL（开源ESP32-CAM球形自平衡机器人）",0.92,6.48,5.35,0.25,10,MUTED,align=PP_ALIGN.CENTER)
specs(s,[
    ("机理", "透明球壳内置自平衡小车/摆杆机构，通过内部重心与轮组驱动球壳滚动。"),
    ("代表", "ESP-ROLL为开源DIY平台；可使用100 mm透明球壳，整机尺寸取决于球壳。"),
    ("速度", "教程未报告标定最高速度，主要面向遥控展示和FPV滚动。"),
    ("控制", "ESP32-CAM、Wi-Fi手机控制、3D打印结构与Arduino固件。"),
    ("精度", "无绝对定位与动态轨迹RMSE；开环滚动和球壳接触难以保证毫米级路径。"),
    ("场景判断", "开源、可改造、外形安全，但低速漂移、滚动惯性和换向滞后仍是核心风险。"),
])
add_text(s,"筛选：适合开源复现和结构探索，不适合作为高精度机械臂替代平台。",7.08,6.34,5.10,0.58,14,RED,True)
footer(s,7,[("ESP-ROLL",URL["esp_roll"]),("模型/照片",URL["esp_roll_files"]),("视频",URL["esp_roll_video"])])

# 8. Special holonomic
s = prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
title(s,"全向球与特殊全向驱动", "连续接触的全向球驱动最接近机械臂平面末端的运动效果")
picture_contain(s, ASSETS/"cellulo_body_epfl.jpg",0.82,1.84,5.55,4.60,URL["cellulo_photo"])
add_text(s,"代表实物：Cellulo（全向球驱动科研平台）",0.92,6.48,5.35,0.25,10,MUTED,align=PP_ALIGN.CENTER)
specs(s,[
    ("机理", "三个全向球驱动单元以120°布置，连续接触地面并独立控制vx、vy、ω。"),
    ("代表", "Cellulo约75 mm、167.8 g；另有双自由度环面轮等实验室特殊机构。"),
    ("速度", "Cellulo实测最大线速度约185 mm/s，适合康复桌面速度范围。"),
    ("控制", "约93 Hz板载定位/控制、Bluetooth、位置与路径跟踪、辅助回驱。"),
    ("精度", "纸面定位约0.27 mm、1.5°；这是感知精度，公开资料未给动态路径RMSE。"),
    ("场景判断", "真全向、连续接触、尺寸可掌握，已有康复研究；主要问题是未商业化。"),
])
add_text(s,"筛选：当前机理上最符合要求，是后续优先获取和验证的方向。",7.08,6.34,5.10,0.58,14,GREEN,True)
footer(s,8,[("Cellulo论文",URL["cellulo"]),("运动控制",URL["cellulo_control"]),("Modulo",URL["modulo"]),("环面轮",URL["torus"])])

# 9. Mechanism comparison
s = prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
title(s,"机理层面的横向筛选", "先判断结构上能否达到目标，再考虑具体产品或科研平台")
headers=["机理","可横移","低速平滑潜力","小型化证据","精度证据","场景结论"]
widths=[2.05,1.28,2.20,2.08,2.32,2.05]
x0=0.72; y0=1.83; x=x0
for h,w in zip(headers,widths):
    add_text(s,h,x,y0,w,0.34,11.7,NAVY,True,PP_ALIGN.CENTER)
    x+=w
rows=[
    ("双轮差速","否","较好","32 mm商品","无机械臂级实证","原型备选"),
    ("三轮全向轮","是","中等","100 mm科研平台","定位有、轨迹不足","第二候选"),
    ("四轮麦克纳姆","是","较差","约200 mm以上","横移轮滑明显","排除"),
    ("球形滚动","非严格","较差","42–60 mm","无精度保证","排除"),
    ("全向球/特殊","是","最好","75 mm科研平台","0.27 mm感知；路径待测","优先"),
]
for i,row in enumerate(rows):
    y=y0+0.54+i*0.90
    ln=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x0),Inches(y-0.13),Inches(sum(widths)),Inches(0.012))
    ln.fill.solid(); ln.fill.fore_color.rgb=LINE; ln.line.fill.background()
    x=x0
    for j,(val,w) in enumerate(zip(row,widths)):
        c=TEXT
        if j==1: c=GREEN if val=="是" else RED if val=="否" else ORANGE
        if j==5: c=GREEN if val=="优先" else ORANGE if "候选" in val else RED
        add_text(s,val,x+0.04,y,w-0.08,0.52,12.4,c,j in (0,5),PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
        x+=w
add_text(s,"结论并不是“全向轮一定最好”，而是连续接触的全向球比离散滚子更有机会实现低速丝滑。",0.90,6.60,11.5,0.36,15,NAVY,True,PP_ALIGN.CENTER)
footer(s,9)

# 10. Selection
s = prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
title(s,"面向康复场景的筛选结果", "按机械臂等效程度、尺寸、安全性和现有证据进行取舍")
sections=[
    ("优先方向", "全向球/特殊全向驱动", "Cellulo类机构同时满足真全向、连续接触、75 mm级尺寸和高频绝对定位。", GREEN),
    ("第二方向", "三轮全向轮", "运动学满足，但现有平台多在100 mm以上；需要重点验证滚子振动、侧滑与低速轨迹。", ORANGE),
    ("快速验证", "双轮差速", "toio尺寸和控制最成熟，可验证康复交互流程，但不能作为机械臂运动等效证据。", BLUE),
    ("不建议", "麦克纳姆 / 球形滚动", "前者过大且低速振动，后者精度和换向不可控，均不适合最终目标。", RED),
]
y=1.92
for head,name,body,color in sections:
    add_text(s,head,0.92,y,1.18,0.32,13,color,True)
    add_text(s,name,2.28,y,2.70,0.34,18,NAVY,True)
    add_text(s,body,5.15,y,7.05,0.62,14.3,TEXT)
    y+=1.14
add_text(s,"当前没有可直接购买、尺寸小于80 mm、又具有机械臂级动态轨迹精度的完整产品。",0.92,6.62,11.35,0.36,16,NAVY,True,PP_ALIGN.CENTER)
footer(s,10)

# 11. Final recommendation
s = prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
add_text(s,"最终建议",0.86,0.76,2.2,0.35,14,TEAL,True)
add_text(s,"先选机理，再选平台",0.84,1.36,7.2,0.72,34,NAVY,True)
add_text(s,"最接近机械臂二维末端效果的不是普通小车，而是全向球/连续接触式全向机构。",0.86,2.34,10.9,0.55,20,TEXT,True)
points=[
    ("01", "优先联系Cellulo/Modulo相关团队，确认平台获取、控制接口和动态轨迹数据。"),
    ("02", "若无法获得，三轮全向轮是可实施的结构备选，但必须做到更小并开展低速滑移测试。"),
    ("03", "toio可用于近期软件与康复流程验证，但汇报中应明确其非全向，不等效于机械臂。"),
]
y=3.40
for n,t in points:
    add_text(s,n,0.92,y,0.55,0.34,14,TEAL,True)
    add_text(s,t,1.66,y,10.45,0.54,16,TEXT)
    y+=0.94
line=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0.92),Inches(6.40),Inches(11.3),Inches(0.035))
line.fill.solid(); line.fill.fore_color.rgb=TEAL; line.line.fill.background()
add_text(s,"需要补齐的关键指标：最低稳定速度、90°换向响应、动态轨迹RMSE、载荷与接触扰动。",0.94,6.64,11.25,0.36,15,NAVY,True)
footer(s,11,[("Cellulo",URL["cellulo"]),("Misaka",URL["misaka"]),("toio",URL["toio"])])

prs.core_properties.title = "康复桌面抓取目标微型移动机器人运动机理分类调研"
prs.core_properties.subject = "差速、全向轮、麦克纳姆、球形与全向球驱动筛选"
prs.core_properties.author = "石慧峰"
prs.save(OUT)
print(OUT)
print("slides", len(prs.slides))
