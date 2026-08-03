from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from PIL import Image

BASE = Path(__file__).resolve().parent
ASSETS = BASE / "assets"
OUT = BASE / "康复桌面抓取目标微型移动机器人调研_汇报版_v1.pptx"

# 16:9 blank deck
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# Visual system
NAVY = RGBColor(23, 50, 77)
TEAL = RGBColor(19, 151, 147)
TEAL_DARK = RGBColor(12, 111, 109)
BLUE = RGBColor(49, 105, 163)
TEXT = RGBColor(31, 43, 55)
MUTED = RGBColor(92, 107, 121)
LIGHT = RGBColor(247, 249, 251)
PALE = RGBColor(232, 241, 245)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(213, 221, 228)
GREEN = RGBColor(46, 125, 85)
ORANGE = RGBColor(225, 137, 45)
RED = RGBColor(190, 69, 76)
FONT = "Microsoft YaHei"

SOURCE_URLS = {
    "ur10": "https://www.universal-robots.com/media/50895/ur10_en.pdf",
    "cellulo": "https://github.com/chili-epfl/cellulo",
    "cellulo_paper": "https://pearl.plymouth.ac.uk/secam-research/1570/",
    "modulo": "https://infoscience.epfl.ch/record/295712/files/IROS2022_Modular_Cellulo.pdf",
    "toio": "https://toio.io/platform/cube/",
    "toio_spec": "https://toio.github.io/toio-spec/en/docs/hardware_other/",
    "toio_motor": "https://toio.github.io/toio-spec/en/docs/ble_motor/",
    "toio_python": "https://pypi.org/project/toio.py/",
    "toio_eval": "https://ramarko.com/portfolio/NetworkBots.pdf",
    "epuck": "https://www.gctronic.com/doc/index.php/e-puck2",
    "epuck_flyer": "https://projects.gctronic.com/epuck2/e-puck2-flyer.pdf",
    "epuck_ros2": "https://github.com/cyberbotics/epuck_ros2",
    "sphero": "https://sphero.com/products/sphero-mini",
    "sphero_chart": "https://dmmedia.sphero.com/email-marketing/Sphero/Robot%20Comparison%20Chart.pdf",
    "sphero_accuracy": "https://help.sphero.com/sphero-support/bolt1-speed-and-distance-accuracy",
    "misaka": "https://arxiv.org/abs/2404.17125",
    "misaka_github": "https://github.com/TingliangZhang/Misaka",
    "zooids": "https://shape.stanford.edu/research/swarm/",
    "zooids_paper": "https://shape.stanford.edu/research/swarm/SwarmUIs_no_copyright.pdf",
}


def add_text(slide, text, x, y, w, h, size=18, color=TEXT, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font=FONT,
             margin=0.02, hyperlink=None, line_spacing=1.0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    if hyperlink:
        r.hyperlink.address = hyperlink
    return box


def add_runs(slide, runs, x, y, w, h, size=12, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    for item in runs:
        r = p.add_run(); r.text = item[0]
        r.font.name = FONT; r.font.size = Pt(item[1] if len(item) > 1 else size)
        r.font.bold = item[2] if len(item) > 2 else False
        r.font.color.rgb = item[3] if len(item) > 3 else MUTED
        if len(item) > 4 and item[4]: r.hyperlink.address = item[4]
    return box


def add_rect(slide, x, y, w, h, fill=WHITE, line=LINE, radius=True, transparency=0):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    try: shp.fill.transparency = transparency
    except Exception: pass
    shp.line.color.rgb = line
    shp.line.width = Pt(0.7)
    return shp


def add_line(slide, x1, y1, x2, y2, color=LINE, width=1.0):
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x1), Inches(y1), Inches(x2-x1), Inches(y2-y1 if y2>y1 else 0.01))
    ln.fill.solid(); ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    return ln


def add_title(slide, title, subtitle=None, section=None):
    if section:
        add_text(slide, section.upper(), 0.78, 0.42, 2.6, 0.3, 10, TEAL, True)
    add_text(slide, title, 0.78, 0.73, 11.8, 0.55, 28, NAVY, True)
    if subtitle:
        add_text(slide, subtitle, 0.80, 1.29, 11.6, 0.38, 12.5, MUTED)
    # top accent
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(0.0), Inches(13.333), Inches(0.10))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()


def add_footer(slide, page, source_items=None):
    y = 7.16
    add_text(slide, f"康复桌面抓取目标微型移动机器人调研  |  {page:02d}", 0.78, y, 4.0, 0.20, 8.5, MUTED)
    if source_items:
        x = 5.0
        for i, (label, url) in enumerate(source_items):
            if i:
                add_text(slide, "·", x, y, 0.12, 0.20, 8.5, MUTED)
                x += 0.16
            width = max(0.75, min(2.1, 0.085 * len(label) + 0.30))
            add_text(slide, label, x, y, width, 0.20, 8.5, BLUE, False, hyperlink=url)
            x += width + 0.08


def add_picture_contain(slide, image_path, x, y, w, h, link=None, bg=WHITE):
    image_path = str(image_path)
    with Image.open(image_path) as im:
        iw, ih = im.size
    add_rect(slide, x, y, w, h, fill=bg, line=LINE, radius=True)
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    pic = slide.shapes.add_picture(image_path, Inches(x + (w-pw)/2), Inches(y + (h-ph)/2), Inches(pw), Inches(ph))
    if link: pic.click_action.hyperlink.address = link
    return pic


def add_picture_cover(slide, image_path, x, y, w, h, link=None):
    image_path = str(image_path)
    with Image.open(image_path) as im:
        iw, ih = im.size
    ratio = w / h
    iratio = iw / ih
    if iratio > ratio:
        ph = h; pw = h * iratio
    else:
        pw = w; ph = w / iratio
    pic = slide.shapes.add_picture(image_path, Inches(x - (pw-w)/2), Inches(y - (ph-h)/2), Inches(pw), Inches(ph))
    # Crop to desired box
    if pw > w:
        crop = (pw - w) / (2 * pw); pic.crop_left = crop; pic.crop_right = crop
    if ph > h:
        crop = (ph - h) / (2 * ph); pic.crop_top = crop; pic.crop_bottom = crop
    if link: pic.click_action.hyperlink.address = link
    return pic


def add_tag(slide, text, x, y, w, color=TEAL, fill=PALE):
    shp = add_rect(slide, x, y, w, 0.33, fill=fill, line=fill, radius=True)
    add_text(slide, text, x, y+0.015, w, 0.27, 10.5, color, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    return shp


def add_bullet_list(slide, items, x, y, w, h, size=15.5, color=TEXT, gap=0.48, bullet_color=TEAL):
    cy = y
    for item in items:
        slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(cy+0.14), Inches(0.09), Inches(0.09)).fill.solid()
        dot = slide.shapes[-1]
        dot.fill.fore_color.rgb = bullet_color; dot.line.fill.background()
        add_text(slide, item, x+0.18, cy, w-0.18, gap, size, color)
        cy += gap


def add_param_panel(slide, rows, x=7.0, y=1.78, w=5.55, h=4.45, verdict=None, verdict_color=TEAL):
    add_rect(slide, x, y, w, h, fill=WHITE, line=LINE, radius=True)
    row_h = (h-0.65) / len(rows)
    for i, (label, value) in enumerate(rows):
        yy = y + 0.18 + i*row_h
        add_text(slide, label, x+0.22, yy, 1.12, row_h-0.05, 11.3, MUTED, True, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, value, x+1.36, yy, w-1.58, row_h-0.05, 13.0, TEXT, False, valign=MSO_ANCHOR.MIDDLE)
        if i < len(rows)-1:
            ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x+0.20), Inches(yy+row_h-0.03), Inches(w-0.40), Inches(0.008))
            ln.fill.solid(); ln.fill.fore_color.rgb = LINE; ln.line.fill.background()
    if verdict:
        band = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x+0.18), Inches(y+h-0.48), Inches(w-0.36), Inches(0.34))
        band.fill.solid(); band.fill.fore_color.rgb = verdict_color; band.line.fill.background()
        add_text(slide, verdict, x+0.24, y+h-0.445, w-0.48, 0.25, 11.5, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)


def add_candidate_slide(title, subtitle, images, rows, verdict, verdict_color, page, tags, sources, image_captions=None):
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = LIGHT
    add_title(slide, title, subtitle, "候选平台")
    if len(images) == 1:
        add_picture_contain(slide, ASSETS/images[0][0], 0.78, 1.78, 5.85, 4.86, images[0][1], bg=WHITE)
        if image_captions:
            add_text(slide, image_captions[0], 0.98, 6.32, 5.45, 0.23, 9, MUTED, align=PP_ALIGN.CENTER)
    else:
        add_picture_contain(slide, ASSETS/images[0][0], 0.78, 1.78, 5.85, 2.68, images[0][1], bg=WHITE)
        add_picture_contain(slide, ASSETS/images[1][0], 0.78, 4.58, 5.85, 2.06, images[1][1], bg=WHITE)
        if image_captions:
            add_text(slide, image_captions[0], 0.98, 4.20, 5.45, 0.20, 8.5, MUTED, align=PP_ALIGN.CENTER)
            add_text(slide, image_captions[1], 0.98, 6.38, 5.45, 0.20, 8.5, MUTED, align=PP_ALIGN.CENTER)
    tx = 7.0
    for tag, width, color, fill in tags:
        add_tag(slide, tag, tx, 1.52, width, color, fill)
        tx += width + 0.10
    add_param_panel(slide, rows, verdict=verdict, verdict_color=verdict_color)
    add_footer(slide, page, sources)
    return slide


# Slide 1 — cover
s = prs.slides.add_slide(blank)
s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY
# Minimal geometry
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(0.0), Inches(0.16), Inches(7.5))
accent.fill.solid(); accent.fill.fore_color.rgb = TEAL; accent.line.fill.background()
add_text(s, "硬件系统调研", 0.95, 1.05, 5.0, 0.45, 15, TEAL, True)
add_text(s, "康复桌面抓取目标\n微型移动机器人", 0.92, 1.58, 7.4, 1.75, 38, WHITE, True)
add_text(s, "从工业机械臂替代需求出发，比较尺寸、全向运动、速度、控制与精度", 0.95, 3.62, 6.4, 0.75, 17, RGBColor(205,220,231))
add_text(s, "汇报人：石慧峰    2026.07", 0.95, 6.45, 5.2, 0.35, 13, RGBColor(205,220,231))
# Three small product visuals on right
add_picture_contain(s, ASSETS/"cellulo_photo2.jpg", 8.10, 0.92, 4.35, 1.82, SOURCE_URLS["cellulo"], bg=RGBColor(15,34,52))
add_picture_contain(s, ASSETS/"toio_cube.png", 8.10, 2.92, 2.05, 2.20, SOURCE_URLS["toio"], bg=WHITE)
add_picture_contain(s, ASSETS/"sphero_mini.png", 10.35, 2.92, 2.10, 2.20, SOURCE_URLS["sphero"], bg=WHITE)
add_text(s, "核心问题", 8.10, 5.42, 1.2, 0.28, 11, TEAL, True)
add_text(s, "有没有足够小、真正全向、又能达到机械臂轨迹效果的桌面平台？", 8.10, 5.78, 4.35, 0.92, 17, WHITE, True)

# Slide 2 — task dimensions
s = prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb = LIGHT
add_title(s, "调研任务与评价维度", "老师提出的五项要求，最终统一落到“能否复现机械臂末端运动效果”", "调研框架")
dims = [
    ("01", "物理尺寸", "目标可被患者自然触碰或抓取，优先3–8 cm级"),
    ("02", "运动方式", "能否真正全向，换向是否需要先调整车体朝向"),
    ("03", "移动速度", "最高速度之外，更关注低速稳定性与加减速平滑性"),
    ("04", "控制方式", "Python/ROS/SDK、无线延迟以及20 Hz RL控制接入"),
    ("05", "运动精度", "严格区分定位分辨率、定位精度和动态轨迹误差"),
]
for i,(num,head,body) in enumerate(dims):
    y=1.78+i*0.96
    add_text(s,num,0.88,y,0.52,0.42,15,TEAL,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
    add_text(s,head,1.55,y,1.65,0.42,17,NAVY,True,valign=MSO_ANCHOR.MIDDLE)
    add_text(s,body,3.18,y,8.85,0.42,15,TEXT,valign=MSO_ANCHOR.MIDDLE)
    if i<4:
        ln=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(1.50),Inches(y+0.58),Inches(10.55),Inches(0.008))
        ln.fill.solid(); ln.fill.fore_color.rgb=LINE; ln.line.fill.background()
add_rect(s,0.88,6.56,11.45,0.40,fill=PALE,line=PALE,radius=True)
add_text(s,"最终判据：任意二维方向连续运动 + 低速丝滑 + 毫米级闭环 + 可接受尺寸",1.05,6.62,11.05,0.27,14,TEAL_DARK,True,PP_ALIGN.CENTER)
add_footer(s,2)

# Slide 3 — UR10 baseline
s = prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb = LIGHT
add_title(s,"机械臂基准：UR10能提供什么", "移动平台不是只要“能到达”，还要复现机械臂末端的二维运动特性", "基准")
add_picture_contain(s, ASSETS/"ur_series.png",0.78,1.72,5.35,4.72,SOURCE_URLS["ur10"],bg=WHITE)
# benchmark numbers
add_text(s,"±0.1 mm",6.65,1.92,2.65,0.72,32,NAVY,True)
add_text(s,"原UR10官方位姿重复精度",6.68,2.58,3.00,0.32,12,MUTED)
add_text(s,"≈ 1 m/s",9.78,1.92,2.20,0.72,32,NAVY,True)
add_text(s,"典型TCP速度",9.80,2.58,2.20,0.32,12,MUTED)
add_rect(s,6.65,3.18,5.65,2.54,fill=WHITE,line=LINE,radius=True)
add_bullet_list(s,[
    "末端可沿任意x–y方向移动，不需要先改变自身朝向",
    "直线、圆弧和复杂轨迹连续，速度与加速度可规划",
    "位置反馈与驱动闭环完整，重复性远优于轮式平台",
    "桌面机器人若存在打滑或非完整约束，视觉效果会明显不同",
],6.92,3.44,5.08,2.05,14.1,gap=0.48)
add_rect(s,0.90,6.45,11.35,0.47,fill=NAVY,line=NAVY,radius=True)
add_text(s,"“达到机械臂效果”至少要求：真正全向运动 + 平滑低速 + 毫米级动态轨迹控制",1.10,6.53,10.95,0.28,14,WHITE,True,PP_ALIGN.CENTER)
add_footer(s,3,[("UR10官方规格",SOURCE_URLS["ur10"])])

# Slide 4 — technology landscape
s = prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb = LIGHT
add_title(s,"三类运动机构的本质差异", "“能向任意方向到达”不等于“能够像机械臂一样全向移动”", "运动学")
cols=[
    ("真正全向", "Cellulo / Modulo\nMisaka / 三全向轮", "可独立控制 vx、vy、ω\n无需先转车头", GREEN, "最接近机械臂"),
    ("双轮差速", "Sony toio\ne-puck2 / Zooids", "只能沿车头方向运动\n横向目标需先转向", ORANGE, "轨迹会出现转向过程"),
    ("球形滚动", "Sphero Mini / BOLT", "外观看似任意方向\n内部转向与球壳惯性明显", RED, "低速与精度不稳定"),
]
for i,(head,names,body,color,verdict) in enumerate(cols):
    x=0.80+i*4.17
    add_rect(s,x,1.82,3.78,4.78,fill=WHITE,line=LINE,radius=True)
    add_text(s,head,x+0.25,2.10,3.28,0.42,21,color,True,PP_ALIGN.CENTER)
    # simple directional graphic
    cx=x+1.89; cy=3.02
    circle=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(cx-0.37),Inches(cy-0.37),Inches(0.74),Inches(0.74))
    circle.fill.solid(); circle.fill.fore_color.rgb=PALE; circle.line.color.rgb=color
    if i==0:
        dirs=[(-0.75,0),(0.75,0),(0,-0.75),(0,0.75),(-0.53,-0.53),(0.53,-0.53),(-0.53,0.53),(0.53,0.53)]
    elif i==1:
        dirs=[(0,-0.78),(0,0.78)]
    else:
        dirs=[(-0.62,-0.45),(0.62,-0.45),(-0.62,0.45),(0.62,0.45)]
    for dx,dy in dirs:
        dot=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(cx+dx-0.045),Inches(cy+dy-0.045),Inches(0.09),Inches(0.09))
        dot.fill.solid(); dot.fill.fore_color.rgb=color; dot.line.fill.background()
    add_text(s,names,x+0.28,3.76,3.22,0.82,16,NAVY,True,PP_ALIGN.CENTER)
    add_text(s,body,x+0.30,4.62,3.18,0.78,13.5,MUTED,False,PP_ALIGN.CENTER)
    add_rect(s,x+0.30,5.72,3.18,0.46,fill=color,line=color,radius=True)
    add_text(s,verdict,x+0.38,5.80,3.02,0.27,12,WHITE,True,PP_ALIGN.CENTER)
add_footer(s,4)

# Slide 5 — Cellulo / Modulo
add_candidate_slide(
    "Cellulo / Modulo Cellulo", "唯一同时具备手持尺寸、真正全向球驱动和高频绝对定位的候选",
    [("cellulo_photo1.jpg",SOURCE_URLS["cellulo"]),("cellulo_v2_exploded.png",SOURCE_URLS["modulo"])],
    [
        ("尺寸/重量","原Cellulo约75 mm宽、30 mm高、约168 g；Modulo约200 g"),
        ("运动方式","三组全向球驱动，vx、vy、ω独立控制；真正3-DOF全向"),
        ("移动速度","约0.2 m/s；Modulo控制实验覆盖50–150 mm/s"),
        ("控制方式","Bluetooth SPP；Unity SDK；论文记录ROS接口与板载闭环"),
        ("定位精度","原版约0.27 mm/1.6°；Modulo约0.17 mm/1.5°，约93 Hz"),
        ("关键限制","公开值主要是定位精度；毫米单位动态轨迹RMSE未充分报告"),
    ],
    "结论：运动学和定位层面最接近机械臂，但当前难以直接采购",GREEN,5,
    [("真全向",0.92,GREEN,RGBColor(229,242,234)),("亚毫米定位",1.18,TEAL_DARK,PALE),("研究平台",0.98,ORANGE,RGBColor(251,240,222))],
    [("Cellulo GitHub",SOURCE_URLS["cellulo"]),("HRI论文",SOURCE_URLS["cellulo_paper"]),("Modulo论文",SOURCE_URLS["modulo"])],
    ["Cellulo用于桌面交互的公开项目照片","Cellulo V2模块化结构图"]
)

# Slide 6 — toio
add_candidate_slide(
    "Sony toio Core Cube", "最小、最易采购、Python生态最完整，但双轮差速无法复现机械臂横向移动",
    [("toio_cube.png",SOURCE_URLS["toio"])],
    [
        ("尺寸/重量","32 × 32 × 20 mm；约30 g；标称最大载荷200 g"),
        ("运动方式","双轮差速；不能直接横移，改变速度方向通常伴随车体转向"),
        ("移动速度","最大直线速度350 mm/s；最低有效电机指令为8"),
        ("控制方式","BLE 4.2；官方toio.py；左右轮、目标坐标和多目标点控制"),
        ("精度信息","检测分辨率约1.42 mm；实证动态位置差约3.8 mm"),
        ("完成条件","官方目标控制：X/Y各在15坐标单位内、角度误差4°内即完成"),
    ],
    "结论：最现实的现成验证平台，但不是真正全向，也不是机械臂级精度",ORANGE,6,
    [("可直接采购",1.10,GREEN,RGBColor(229,242,234)),("非全向",0.82,RED,RGBColor(250,231,233)),("官方Python",1.05,TEAL_DARK,PALE)],
    [("产品页",SOURCE_URLS["toio"]),("硬件规格",SOURCE_URLS["toio_spec"]),("电机协议",SOURCE_URLS["toio_motor"]),("Python库",SOURCE_URLS["toio_python"])],
    ["Sony官方Core Cube产品图"]
)

# Slide 7 — e-puck2
add_candidate_slide(
    "GCtronic e-puck2", "成熟的微型科研平台，接口开放，但本质仍是非全向差速机器人",
    [("epuck2.png",SOURCE_URLS["epuck"])],
    [
        ("尺寸/重量","直径70 mm、高45 mm；约130 g"),
        ("运动方式","双轮差速；两台步进电机，不能保持朝向进行横向平移"),
        ("移动速度","最大约154 mm/s；软件最小命令可到1 step/s"),
        ("控制方式","USB、Bluetooth、BLE、Wi-Fi；C/C++、Python、ROS 2、Webots"),
        ("分辨率","约0.128–0.13 mm/电机步，只是理论轮端位移，不是定位精度"),
        ("关键限制","无内置绝对定位；步数计数不能检测打滑或实际轮胎位移"),
    ],
    "结论：科研控制生态优秀，但运动方式和精度闭环都达不到机械臂效果",ORANGE,7,
    [("科研成熟",0.92,GREEN,RGBColor(229,242,234)),("ROS 2",0.72,TEAL_DARK,PALE),("非全向",0.82,RED,RGBColor(250,231,233))],
    [("官方文档",SOURCE_URLS["epuck"]),("官方规格",SOURCE_URLS["epuck_flyer"]),("ROS 2驱动",SOURCE_URLS["epuck_ros2"])],
    ["GCtronic官方e-puck2结构与功能图"]
)

# Slide 8 — Sphero
add_candidate_slide(
    "Sphero Mini / BOLT", "球形外观和抓握尺寸理想，但厂商不提供精密位置或轨迹指标",
    [("sphero_life.jpg",SOURCE_URLS["sphero"]),("sphero_mini.png",SOURCE_URLS["sphero"])],
    [
        ("尺寸/重量","Mini直径42 mm、46 g；BOLT直径73 mm、约200 g"),
        ("运动方式","内部驱动带动球壳滚动；表观任意方向，但不是独立3-DOF全向"),
        ("移动速度","Mini最大约1 m/s；BOLT标称最大约2 m/s"),
        ("控制方式","BLE；Sphero Edu；Blocks/JavaScript/Python教学控制"),
        ("精度信息","无Mini位置精度指标；BOLT官方明确不是精密测量仪器"),
        ("误差来源","桌面材质、坡度、电量、温度、程序时序和电机差异"),
    ],
    "结论：尺寸和安全外形很好，但低速、换向和轨迹精度不可保证",RED,8,
    [("尺寸理想",0.92,GREEN,RGBColor(229,242,234)),("球形滚动",1.02,ORANGE,RGBColor(251,240,222)),("无精度保证",1.12,RED,RGBColor(250,231,233))],
    [("Mini产品页",SOURCE_URLS["sphero"]),("规格对比",SOURCE_URLS["sphero_chart"]),("精度声明",SOURCE_URLS["sphero_accuracy"])],
    ["Sphero Mini桌面使用场景","Sphero Mini官方产品图"]
)

# Slide 9 — Misaka
add_candidate_slide(
    "Misaka开源全向平台", "架构方向正确：三全向轮 + 纸面微点阵定位，但尺寸和验证成熟度不足",
    [("misaka_img_6.jpeg",SOURCE_URLS["misaka"]),("misaka_img_12.jpeg",SOURCE_URLS["misaka"])],
    [
        ("尺寸/重量","直径100 mm、高50 mm；论文未报告整机重量"),
        ("运动方式","三只38 mm全向轮；理论上可独立控制vx、vy、ω"),
        ("移动速度","最大约200 mm/s；连续运动续航约30 min"),
        ("控制方式","ATmega2560/Arduino、XBee；可扩展ESP32；Jupyter代码开源"),
        ("定位方式","SONiX OID传感器读取纸面微点阵，输出绝对位置"),
        ("精度限制","论文仅称高精度，没有给出Misaka整机自身的实测定位或轨迹误差"),
    ],
    "结论：可作为开源全向技术参考，但100 mm偏大，精度证据不足",ORANGE,9,
    [("真全向",0.92,GREEN,RGBColor(229,242,234)),("开源硬件",1.00,TEAL_DARK,PALE),("尺寸偏大",0.92,RED,RGBColor(250,231,233))],
    [("arXiv论文",SOURCE_URLS["misaka"]),("项目GitHub",SOURCE_URLS["misaka_github"])],
    ["Misaka机器人主控与全向底盘实物","Misaka主要硬件模块"]
)

# Slide 10 — Zooids
add_candidate_slide(
    "Zooids微型桌面机器人", "体积极小、速度很高，但依赖定制投影定位系统且仍为双轮差速",
    [("zooids_teaser.png",SOURCE_URLS["zooids"]),("zooids_exploded.png",SOURCE_URLS["zooids"])],
    [
        ("尺寸/重量","直径26 mm、高21 mm；约12 g"),
        ("运动方式","双轮差速；不是全向，论文也将非完整约束列为局限"),
        ("移动速度","最大约740 mm/s；平滑可控的常用速度约440 mm/s"),
        ("控制方式","中央控制器 + 无线通信 + 高速DLP结构光投影跟踪"),
        ("跟踪信息","投影/跟踪空间分辨率约1.1 mm；不等于实际运动轨迹误差"),
        ("关键限制","研究原型、非商品；基础设施复杂，难以直接移植到康复系统"),
    ],
    "结论：尺寸最小，但非全向且系统复杂，不适合作为直接替代平台",RED,10,
    [("超小尺寸",0.98,GREEN,RGBColor(229,242,234)),("非全向",0.82,RED,RGBColor(250,231,233)),("研究原型",0.98,ORANGE,RGBColor(251,240,222))],
    [("项目主页",SOURCE_URLS["zooids"]),("UIST论文",SOURCE_URLS["zooids_paper"])],
    ["Zooids群体桌面交互演示","Zooids微型机器人内部结构"]
)

# Slide 11 — comparison table
s = prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb = LIGHT
add_title(s,"候选平台横向对比", "最关键的矛盾是：现成商用品不全向，真正全向的平台尚未成熟商业化", "综合比较")
columns=[("平台",0.95),("尺寸",1.30),("运动方式",1.72),("速度",1.22),("控制",2.10),("精度证据",2.75),("判断",1.72)]
x0=0.68; y0=1.72; row_h=0.72
# header
x=x0
for label,w in columns:
    add_rect(s,x,y0,w,row_h,fill=NAVY,line=WHITE,radius=False)
    add_text(s,label,x+0.04,y0+0.15,w-0.08,0.30,11.5,WHITE,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
    x+=w
rows=[
    ("Cellulo","75 mm","真全向球驱动","0.20 m/s","BT / Unity / ROS","0.17–0.27 mm定位；动态轨迹未证实","最接近"),
    ("toio","32 mm","双轮差速","0.35 m/s","BLE / 官方Python","约1.42 mm分辨率；实证约3.8 mm","最现实"),
    ("e-puck2","70 mm","双轮差速","0.154 m/s","多通信 / ROS 2","0.13 mm/步不是实际精度","不匹配"),
    ("Sphero","42/73 mm","球形滚动","1–2 m/s","BLE / Edu","厂商无位置精度保证","不推荐"),
    ("Misaka","100 mm","三全向轮","0.20 m/s","Arduino / XBee","无整机实测精度","参考"),
    ("Zooids","26 mm","双轮差速","0.74 m/s","投影跟踪 / 定制","约1.1 mm系统分辨率","不直接用"),
]
for ri,row in enumerate(rows):
    y=y0+row_h*(ri+1); x=x0
    bg=WHITE if ri%2==0 else RGBColor(240,245,248)
    for ci,((label,w),val) in enumerate(zip(columns,row)):
        add_rect(s,x,y,w,row_h,fill=bg,line=LINE,radius=False)
        col=TEXT; bold=(ci==0)
        if ci==6:
            col = GREEN if val in ("最接近","最现实") else ORANGE if val=="参考" else RED
            bold=True
        add_text(s,val,x+0.05,y+0.10,w-0.10,row_h-0.15,10.5,col,bold,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
        x+=w
add_rect(s,0.78,6.90,11.78,0.20,fill=TEAL,line=TEAL,radius=True)
add_footer(s,11)

# Slide 12 — recommendation
s = prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb = LIGHT
add_title(s,"推荐结论", "没有发现同时满足“小尺寸、真全向、机械臂级精度、可直接采购”的现成产品", "决策")
# left technical ideal
add_rect(s,0.78,1.78,5.82,4.95,fill=WHITE,line=LINE,radius=True)
add_text(s,"技术最优",1.05,2.06,1.4,0.35,13,TEAL,True)
add_text(s,"Cellulo / Modulo Cellulo",1.05,2.47,4.95,0.48,24,NAVY,True)
add_bullet_list(s,[
    "真正3-DOF全向，最接近机械臂平面末端",
    "约93 Hz纸面绝对定位，位置测量达到亚毫米级",
    "75 mm手持尺寸仍可用于掌握类康复任务",
    "主要障碍是不可常规采购，动态轨迹精度仍需验证",
],1.08,3.22,5.05,2.20,14.2,gap=0.50)
add_rect(s,1.05,5.62,5.28,0.68,fill=RGBColor(229,242,234),line=RGBColor(229,242,234),radius=True)
add_text(s,"若“机械臂效果”是硬性要求，应优先联系EPFL获取平台或技术合作",1.24,5.73,4.90,0.44,12,GREEN,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
# right realistic
add_rect(s,6.85,1.78,5.70,4.95,fill=WHITE,line=LINE,radius=True)
add_text(s,"现实采购",7.12,2.06,1.4,0.35,13,ORANGE,True)
add_text(s,"Sony toio Core Cube",7.12,2.47,4.85,0.48,24,NAVY,True)
add_bullet_list(s,[
    "32 mm、30 g，是尺寸最合适的现成平台",
    "官方Python、BLE和目标坐标控制便于快速接入",
    "绝对定位优于普通差速小车，已有几毫米级实证",
    "但无法横向平移，不能严格复现机械臂轨迹",
],7.15,3.22,4.92,2.20,14.2,gap=0.50,bullet_color=ORANGE)
add_rect(s,7.12,5.62,5.16,0.68,fill=RGBColor(251,240,222),line=RGBColor(251,240,222),radius=True)
add_text(s,"适合作为概念验证和算法迁移平台，不宜宣称达到机械臂效果",7.30,5.73,4.80,0.44,12,ORANGE,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
add_footer(s,12,[("Cellulo项目",SOURCE_URLS["cellulo"]),("toio产品",SOURCE_URLS["toio"])])

# Slide 13 — test and supplier checklist
s = prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb = LIGHT
add_title(s,"正式选型前必须补齐的数据", "现有厂商资料多数只给最高速度或定位分辨率，不能代替动态运动测试", "验证清单")
items=[
    ("低速稳定性","5 / 10 / 20 mm/s是否连续，无爬行、抖动和明显速度波纹"),
    ("重复定位","同一目标点重复30–100次，报告均值、标准差和最大误差"),
    ("动态轨迹","直线、圆形、S形在50–200 mm/s下的RMSE和最大横向误差"),
    ("瞬时换向","90°与180°改变速度方向时的延迟、停顿、超调和jerk"),
    ("控制链路","无线指令延迟、抖动、丢包率以及20 Hz RL控制的持续稳定性"),
    ("抓取扰动","附加目标外壳、手部接触和载荷变化后的位置与速度误差"),
]
for i,(head,body) in enumerate(items):
    col=i%2; row=i//2
    x=0.82+col*6.18; y=1.80+row*1.52
    add_rect(s,x,y,5.70,1.22,fill=WHITE,line=LINE,radius=True)
    add_text(s,f"{i+1:02d}",x+0.22,y+0.22,0.48,0.34,13,TEAL,True,PP_ALIGN.CENTER)
    add_text(s,head,x+0.78,y+0.17,1.35,0.32,15,NAVY,True)
    add_text(s,body,x+2.08,y+0.15,3.30,0.72,12.5,MUTED)
add_rect(s,0.90,6.55,11.45,0.42,fill=NAVY,line=NAVY,radius=True)
add_text(s,"建议功能等效目标：真全向｜定位≤1 mm｜动态轨迹RMSE≤2 mm｜控制延迟≤50 ms",1.08,6.62,11.05,0.26,13.5,WHITE,True,PP_ALIGN.CENTER)
add_footer(s,13)

# Slide 14 — conclusion and source gateway
s = prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY
add_text(s,"调研结论",0.86,0.70,3.0,0.45,14,TEAL,True)
add_text(s,"当前市场存在明确缺口",0.84,1.32,7.2,0.72,32,WHITE,True)
add_text(s,"没有3–8 cm商业成品同时实现真全向、丝滑低速和机械臂级动态精度。",0.86,2.12,8.40,0.55,18,RGBColor(211,224,233),True)
conclusions=[
    ("01","Cellulo / Modulo最接近机械臂效果，但属于研究平台，采购与动态精度证据是主要缺口。"),
    ("02","Sony toio是最现实的商用选择，适合快速验证，但双轮差速决定了它不能完全替代机械臂。"),
    ("03","Sphero、e-puck2、Zooids和Misaka分别受限于精度、非全向、基础设施或尺寸成熟度。"),
]
for i,(n,t) in enumerate(conclusions):
    y=3.02+i*0.86
    add_text(s,n,0.90,y,0.55,0.38,14,TEAL,True,PP_ALIGN.CENTER)
    add_text(s,t,1.60,y-0.01,9.95,0.56,15.5,WHITE,False)
add_text(s,"主要来源（点击打开）",0.88,6.08,2.05,0.30,11,TEAL,True)
links=[
    ("Cellulo",SOURCE_URLS["cellulo"]),("Modulo",SOURCE_URLS["modulo"]),("Sony toio",SOURCE_URLS["toio"]),
    ("e-puck2",SOURCE_URLS["epuck"]),("Sphero",SOURCE_URLS["sphero"]),("Misaka",SOURCE_URLS["misaka"]),("Zooids",SOURCE_URLS["zooids"])
]
x=0.90
for label,url in links:
    w=0.52+0.095*len(label)
    add_text(s,label,x,6.52,w,0.28,10.5,RGBColor(164,214,222),False,hyperlink=url)
    x+=w+0.22
add_text(s,"谢谢老师，请指导",9.25,6.75,3.20,0.35,15,WHITE,True,PP_ALIGN.RIGHT)

# Document properties
prs.core_properties.title = "康复桌面抓取目标微型移动机器人调研"
prs.core_properties.subject = "工业机械臂替代方案硬件调研"
prs.core_properties.author = "石慧峰"
prs.core_properties.keywords = "康复机器人, 桌面移动机器人, Cellulo, toio, 全向移动"
prs.save(OUT)
print(OUT)
print("slides", len(prs.slides))
