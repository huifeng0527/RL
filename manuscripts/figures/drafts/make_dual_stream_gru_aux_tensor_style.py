from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[3]
OUT_DIR = ROOT / "manuscripts" / "figures" / "paper_ready"
OUT = OUT_DIR / "fig_dual_stream_gru_aux_architecture_tensor_style.pptx"


def I(v: float):
    return Inches(v)


C = {
    "ink": RGBColor(32, 37, 45),
    "muted": RGBColor(95, 106, 123),
    "light": RGBColor(231, 236, 244),
    "blue": RGBColor(38, 97, 145),
    "blue2": RGBColor(202, 225, 244),
    "blue3": RGBColor(237, 246, 252),
    "gold": RGBColor(133, 103, 43),
    "gold2": RGBColor(242, 215, 150),
    "gold3": RGBColor(255, 248, 226),
    "teal": RGBColor(31, 131, 130),
    "teal2": RGBColor(202, 235, 232),
    "teal3": RGBColor(234, 248, 247),
    "red": RGBColor(161, 61, 78),
    "red2": RGBColor(247, 218, 224),
    "red3": RGBColor(255, 242, 245),
    "navy": RGBColor(58, 75, 116),
    "navy2": RGBColor(225, 232, 248),
    "purple": RGBColor(91, 77, 126),
    "purple2": RGBColor(237, 231, 248),
    "green": RGBColor(55, 142, 76),
    "green2": RGBColor(223, 244, 226),
    "white": RGBColor(255, 255, 255),
}


def text(slide, x, y, w, h, s, size=8, color=None, bold=False, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    first = True
    for line in str(s).split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        r = p.add_run()
        r.text = line
        r.font.name = "Arial"
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color or C["ink"]
    return shape


def rect(slide, x, y, w, h, fill=None, line=None, lw=0.9, dashed=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(x), I(y), I(w), I(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line or C["ink"]
    shp.line.width = Pt(lw)
    if dashed:
        shp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return shp


def line(slide, x1, y1, x2, y2, color=None, width=0.9, dashed=False, arrow=False, begin=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
    conn.line.color.rgb = color or C["ink"]
    conn.line.width = Pt(width)
    if dashed:
        conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if arrow or begin:
        ln = conn.line._get_or_add_ln()
        if arrow:
            tail = ln.find(qn("a:tailEnd"))
            if tail is None:
                tail = OxmlElement("a:tailEnd")
                ln.append(tail)
            tail.set("type", "triangle")
        if begin:
            head = ln.find(qn("a:headEnd"))
            if head is None:
                head = OxmlElement("a:headEnd")
                ln.append(head)
            head.set("type", "triangle")
    return conn


def tensor_stack(slide, x, y, w, h, label, dim, color, fill, plates=4, dx=0.055, dy=-0.045, slices=False):
    for i in range(plates - 1, 0, -1):
        rect(slide, x + dx * i, y + dy * i, w, h, fill=RGBColor(250, 253, 255), line=color, lw=0.6)
    top = rect(slide, x, y, w, h, fill=fill, line=color, lw=1.0)
    if slices:
        for j in range(1, 5):
            yy = y + h * j / 5.0
            line(slide, x + 0.04, yy, x + w - 0.04, yy, color=RGBColor(188, 207, 226), width=0.45)
    text(slide, x - 0.18, y + h + 0.06, w + 0.36, 0.22, label, size=7.4, color=color, bold=True)
    text(slide, x - 0.12, y + h + 0.28, w + 0.24, 0.18, dim, size=6.7, color=C["muted"])
    return top


def vector(slide, x, y, w, h, label, dim, color, fill, slots=5):
    shp = rect(slide, x, y, w, h, fill=fill, line=color, lw=1.0)
    for j in range(1, slots):
        yy = y + h * j / slots
        line(slide, x, yy, x + w, yy, color=RGBColor(210, 221, 232), width=0.4)
    text(slide, x - 0.26, y + h + 0.06, w + 0.52, 0.20, label, size=7.3, color=color, bold=True)
    text(slide, x - 0.22, y + h + 0.26, w + 0.44, 0.16, dim, size=6.5, color=C["muted"])
    return shp


def fc_bar(slide, x, y, h, label, dim, color, fill):
    rect(slide, x, y, 0.16, h, fill=fill, line=color, lw=0.85)
    text(slide, x - 0.28, y - 0.18, 0.72, 0.18, label, size=7.0, color=color, bold=True)
    text(slide, x - 0.34, y + h + 0.05, 0.84, 0.16, dim, size=6.2, color=C["muted"])


def gru_cells(slide, x, y, color):
    xs = [x, x + 0.56, x + 1.12]
    labels = ["x_{t-15}", "...", "x_t"]
    for idx, xx in enumerate(xs):
        vector(slide, xx, y, 0.18, 0.62, labels[idx], "", color, C["gold3"], slots=4)
        if idx < 2:
            line(slide, xx + 0.18, y + 0.31, xs[idx + 1] - 0.03, y + 0.31, color=color, width=0.8, arrow=True)
    line(slide, xs[0] + 0.09, y - 0.06, xs[2] + 0.09, y - 0.06, color=color, width=0.65, dashed=True, arrow=True)
    text(slide, x + 0.16, y - 0.35, 0.95, 0.22, "GRU over 16 steps", size=7.2, color=color, bold=True)


def dot(slide, x, y, r, fill, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, I(x), I(y), I(r), I(r))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line_color or fill
    shp.line.width = Pt(0.8)
    return shp


def build():
    prs = Presentation()
    prs.slide_width = I(13.33)
    prs.slide_height = I(7.20)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C["white"]

    text(slide, 0.30, 0.08, 12.75, 0.28, "Dual-stream recurrent policy network with auxiliary future dynamics", 14, C["ink"], True)
    rect(slide, 0.22, 0.52, 12.88, 6.34, fill=None, line=RGBColor(155, 164, 178), lw=1.0, dashed=True)

    # Branch frames, using only dashed boundaries like the reference figure.
    rect(slide, 0.52, 0.88, 4.82, 2.35, fill=None, line=C["blue"], lw=0.9, dashed=True)
    rect(slide, 0.52, 3.55, 4.82, 2.32, fill=None, line=C["gold"], lw=0.9, dashed=True)
    rect(slide, 5.65, 1.95, 2.35, 2.62, fill=None, line=C["teal"], lw=0.9, dashed=True)
    rect(slide, 8.38, 0.88, 4.35, 2.35, fill=None, line=C["red"], lw=0.9, dashed=True)
    rect(slide, 8.38, 3.66, 4.35, 2.22, fill=None, line=C["navy"], lw=0.9, dashed=True)
    text(slide, 1.80, 0.96, 2.10, 0.22, "Spatial branch", 11, C["blue"], True)
    text(slide, 1.33, 3.64, 3.20, 0.22, "Temporal interaction-history branch", 11, C["gold"], True)
    text(slide, 6.02, 2.04, 1.60, 0.20, "Fusion module", 10.5, C["teal"], True)
    text(slide, 9.12, 0.96, 3.00, 0.20, "Auxiliary task module", 10.5, C["red"], True)
    text(slide, 9.35, 3.76, 2.60, 0.20, "PPO actor-critic module", 10.5, C["navy"], True)

    # Observation split.
    text(slide, 0.68, 0.62, 1.25, 0.18, "observation o_t", 7.2, C["muted"], True, PP_ALIGN.LEFT)
    vector(slide, 0.82, 1.42, 0.22, 1.16, "s_t", "12-D", C["blue"], C["blue3"], slots=6)
    tensor_stack(slide, 0.70, 4.24, 0.54, 0.82, "h_{t-15:t}", "16 x C", C["gold"], C["gold3"], plates=5, dx=0.05, dy=-0.05, slices=True)
    line(slide, 0.92, 0.86, 0.92, 1.36, color=RGBColor(174, 184, 196), width=0.55, dashed=True, arrow=True)
    line(slide, 0.92, 0.86, 0.92, 4.16, color=RGBColor(174, 184, 196), width=0.55, dashed=True, arrow=True)
    text(slide, 1.08, 2.52, 0.82, 0.16, "state scalars", 6.4, C["blue"], False, PP_ALIGN.LEFT)
    text(slide, 1.16, 5.15, 1.26, 0.26, "relative interaction\nmotion", 6.3, C["gold"], False, PP_ALIGN.LEFT)

    # Spatial branch, tensors and thin FC bars.
    fc_bar(slide, 2.02, 1.40, 1.08, "FC", "12 -> 32", C["blue"], C["blue2"])
    fc_bar(slide, 2.56, 1.40, 1.08, "LN/ReLU", "", C["blue"], C["blue3"])
    tensor_stack(slide, 3.28, 1.48, 0.46, 0.94, "z_s", "32", C["blue"], C["blue2"], plates=4, dx=0.055, dy=-0.04, slices=True)
    line(slide, 1.05, 2.00, 1.98, 2.00, C["blue"], 1.0, arrow=True)
    line(slide, 2.18, 2.00, 2.52, 2.00, C["blue"], 0.9, arrow=True)
    line(slide, 2.72, 2.00, 3.25, 2.00, C["blue"], 1.0, arrow=True)

    # Temporal branch.
    gru_cells(slide, 1.86, 4.28, C["gold"])
    vector(slide, 3.50, 4.12, 0.22, 1.02, "h_T", "64", C["gold"], C["gold2"], slots=5)
    fc_bar(slide, 4.18, 4.14, 0.98, "FC", "64 -> 32", C["gold"], C["gold2"])
    tensor_stack(slide, 4.70, 4.26, 0.34, 0.74, "z_h", "32", C["gold"], C["gold2"], plates=3, dx=0.045, dy=-0.035, slices=True)
    line(slide, 1.25, 4.65, 1.83, 4.65, C["gold"], 1.0, arrow=True)
    line(slide, 3.03, 4.59, 3.46, 4.59, C["gold"], 1.0, arrow=True)
    line(slide, 3.74, 4.63, 4.14, 4.63, C["gold"], 1.0, arrow=True)
    line(slide, 4.36, 4.63, 4.68, 4.63, C["gold"], 1.0, arrow=True)

    # Fusion module.
    tensor_stack(slide, 5.92, 2.56, 0.32, 1.12, "[z_s; z_h]", "64", C["teal"], C["teal3"], plates=4, dx=0.05, dy=-0.035, slices=True)
    fc_bar(slide, 6.80, 2.50, 1.18, "FC", "64 -> 64", C["teal"], C["teal2"])
    tensor_stack(slide, 7.38, 2.70, 0.36, 0.84, "f_t", "64", C["teal"], C["teal2"], plates=3, dx=0.045, dy=-0.035, slices=True)
    line(slide, 3.76, 1.95, 5.90, 3.02, C["blue"], 1.1, arrow=True)
    line(slide, 5.04, 4.63, 5.91, 3.40, C["gold"], 1.1, arrow=True)
    line(slide, 6.25, 3.12, 6.76, 3.12, C["teal"], 1.0, arrow=True)
    line(slide, 6.98, 3.12, 7.36, 3.12, C["teal"], 1.0, arrow=True)

    # Auxiliary head as thin tensors.
    fc_bar(slide, 8.74, 1.54, 1.12, "FC", "64 -> 64", C["red"], C["red2"])
    fc_bar(slide, 9.54, 1.25, 0.82, "FC", "64 -> H*2", C["red"], C["red2"])
    tensor_stack(slide, 10.22, 1.30, 0.46, 0.74, "D_hat", "H x 2", C["red"], C["red3"], plates=5, dx=0.05, dy=-0.035, slices=True)
    fc_bar(slide, 9.54, 2.36, 0.58, "FC", "", C["red"], C["red3"])
    dot(slide, 10.35, 2.55, 0.13, C["red2"], C["red"])
    text(slide, 10.15, 2.72, 0.52, 0.16, "risk", 6.2, C["red"])
    tensor_stack(slide, 11.45, 1.42, 0.38, 0.64, "targets", "D, y", C["red"], C["white"], plates=3, dx=0.04, dy=-0.03, slices=True)
    text(slide, 8.98, 3.02, 2.70, 0.15, "aux loss: future trajectory + catch-risk", 6.5, C["red"])
    line(slide, 7.73, 3.08, 8.70, 2.10, C["red"], 1.0, arrow=True)
    line(slide, 8.90, 2.10, 9.50, 1.66, C["red"], 0.9, arrow=True)
    line(slide, 8.90, 2.10, 9.50, 2.65, C["red"], 0.9, arrow=True)
    line(slide, 9.70, 1.66, 10.18, 1.66, C["red"], 0.9, arrow=True)
    line(slide, 9.70, 2.65, 10.32, 2.61, C["red"], 0.9, arrow=True)
    line(slide, 11.43, 1.74, 10.75, 1.66, C["red"], 0.75, dashed=True, arrow=True, begin=True)

    # PPO heads.
    fc_bar(slide, 8.74, 4.30, 0.82, "", "", C["navy"], C["navy2"])
    text(slide, 8.42, 4.10, 0.88, 0.16, "Actor head", 6.7, C["navy"], True)
    text(slide, 8.50, 5.14, 0.70, 0.13, "64 -> 2", 5.8, C["muted"])
    rect(slide, 9.60, 4.46, 0.20, 0.52, fill=C["white"], line=C["navy"], lw=0.9)
    line(slide, 9.60, 4.72, 9.80, 4.72, C["light"], 0.4)
    text(slide, 9.43, 5.03, 0.55, 0.14, "a_t", 6.8, C["navy"], True)
    text(slide, 9.32, 5.18, 0.78, 0.14, "2-D action", 5.8, C["muted"])
    fc_bar(slide, 8.74, 5.30, 0.38, "", "", C["purple"], C["purple2"])
    text(slide, 8.42, 5.10, 0.88, 0.16, "Critic head", 6.7, C["purple"], True)
    dot(slide, 9.64, 5.39, 0.13, C["purple2"], C["purple"])
    text(slide, 9.42, 5.56, 0.62, 0.13, "V(o_t)", 6.0, C["purple"])
    tensor_stack(slide, 11.20, 4.48, 0.42, 0.60, "rollout", "o,a,r,V", C["navy"], C["white"], plates=3, dx=0.04, dy=-0.03, slices=True)
    text(slide, 10.70, 5.64, 1.46, 0.15, "PPO update", 6.8, C["navy"], True)
    line(slide, 7.73, 3.12, 8.70, 4.72, C["navy"], 1.0, arrow=True)
    line(slide, 7.73, 3.12, 8.70, 5.49, C["purple"], 1.0, arrow=True)
    line(slide, 8.90, 4.71, 9.56, 4.71, C["navy"], 0.9, arrow=True)
    line(slide, 8.90, 5.49, 9.60, 5.45, C["purple"], 0.9, arrow=True)
    line(slide, 9.82, 4.72, 11.18, 4.72, C["navy"], 0.75, dashed=True, arrow=True)
    line(slide, 11.42, 5.10, 11.42, 5.55, C["navy"], 0.75, dashed=True, arrow=True)
    line(slide, 11.16, 5.72, 8.96, 5.72, C["navy"], 0.75, dashed=True, arrow=True)

    # Gradient hint, deliberately lightweight.
    line(slide, 6.96, 6.42, 6.96, 3.64, C["teal"], 0.65, dashed=True, arrow=True)
    text(slide, 3.60, 6.27, 6.65, 0.24, "L_total = L_PPO + lambda_traj L_traj + lambda_risk L_risk", 8.2, C["muted"])
    text(slide, 0.52, 6.56, 4.30, 0.15, "thin plates denote tensors / hidden features", 6.2, RGBColor(134, 145, 160), False, PP_ALIGN.LEFT)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
