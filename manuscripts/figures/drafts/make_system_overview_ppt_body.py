from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[3]
OUT = ROOT / "manuscripts" / "figures" / "paper_ready" / "fig_overview_system_framework_ppt_body_spatial_temporal_cues.pptx"


def rgb(value):
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


C = {
    "black": rgb("#222222"),
    "gray": rgb("#666a73"),
    "light_gray": rgb("#f4f6f8"),
    "mid_gray": rgb("#c5ccd6"),
    "red": rgb("#c92525"),
    "blue": rgb("#2f6fa7"),
    "blue_light": rgb("#e9f2fb"),
    "orange": rgb("#cf6f2e"),
    "orange_light": rgb("#fff0e4"),
    "green": rgb("#3a8c4a"),
    "green_light": rgb("#e9f5eb"),
    "purple": rgb("#7356b6"),
    "purple_light": rgb("#f0ecfa"),
}


def I(value):
    return Inches(value)


def style_text(shape, text, size=10, color=None, bold=False, align=PP_ALIGN.CENTER):
    color = color or C["black"]
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = I(0.04)
    tf.margin_right = I(0.04)
    tf.margin_top = I(0.02)
    tf.margin_bottom = I(0.02)
    lines = str(text).split("\n")
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        run = p.add_run()
        run.text = line
        run.font.name = "Arial"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def add_text(slide, x, y, w, h, text, size=10, color=None, bold=False, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    style_text(shape, text, size=size, color=color, bold=bold, align=align)
    return shape


def add_box(
    slide,
    x,
    y,
    w,
    h,
    text="",
    fill=None,
    line=None,
    line_width=1.25,
    radius=True,
    dash=False,
    size=10,
    color=None,
    bold=False,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, I(x), I(y), I(w), I(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill or RGBColor(255, 255, 255)
    shape.line.color.rgb = line or C["black"]
    shape.line.width = Pt(line_width)
    if dash:
        shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if text:
        style_text(shape, text, size=size, color=color, bold=bold)
    return shape


def add_panel(slide, x, y, w, h, tag, title):
    add_box(slide, x, y, w, h, fill=RGBColor(255, 255, 255), line=C["black"], line_width=1.6)
    add_text(slide, x + 0.14, y + 0.10, 0.25, 0.22, tag, size=14, bold=True)
    add_text(slide, x + 0.35, y + 0.16, w - 0.7, 0.38, title, size=15.5, color=C["red"], bold=True)


def add_arrow(slide, x, y, w, h, direction, color):
    kind = {
        "right": MSO_SHAPE.RIGHT_ARROW,
        "left": MSO_SHAPE.LEFT_ARROW,
        "down": MSO_SHAPE.DOWN_ARROW,
        "up": MSO_SHAPE.UP_ARROW,
    }[direction]
    shape = slide.shapes.add_shape(kind, I(x), I(y), I(w), I(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    shape.line.width = Pt(0.5)
    return shape


def add_line(slide, x1, y1, x2, y2, color=None, width=0.65, dash=False):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
    line.line.color.rgb = color or C["mid_gray"]
    line.line.width = Pt(width)
    if dash:
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return line


def add_dot(slide, x, y, d, color, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, I(x), I(y), I(d), I(d))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = line or color
    shape.line.width = Pt(0.7)
    return shape


def add_placeholder(slide, x, y, w, h, text):
    shape = add_box(
        slide,
        x,
        y,
        w,
        h,
        text,
        fill=C["light_gray"],
        line=C["mid_gray"],
        line_width=1.1,
        radius=False,
        dash=True,
        size=9.5,
        color=C["gray"],
        bold=True,
    )
    return shape


def small_down_arrow(slide, x, y, color):
    add_arrow(slide, x, y, 0.13, 0.16, "down", color)


def build():
    prs = Presentation()
    prs.slide_width = I(10)
    prs.slide_height = I(8)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Layout grid.
    lx, rx = 0.32, 5.12
    ty, by = 0.48, 4.14
    pw, ph = 4.56, 3.42

    # Top legend.
    add_dot(slide, 0.48, 0.13, 0.10, C["blue"])
    add_text(slide, 0.61, 0.08, 0.92, 0.20, "robot / microrobot", size=8.4, color=C["gray"], align=PP_ALIGN.LEFT)
    add_dot(slide, 1.70, 0.13, 0.10, C["orange"])
    add_text(slide, 1.83, 0.08, 0.76, 0.20, "patient hand", size=8.4, color=C["gray"], align=PP_ALIGN.LEFT)
    add_dot(slide, 2.62, 0.13, 0.10, C["green"])
    add_text(slide, 2.75, 0.08, 0.74, 0.20, "vision / state", size=8.4, color=C["gray"], align=PP_ALIGN.LEFT)

    add_panel(slide, lx, ty, pw, ph, "a", "Simulation Environment")
    add_panel(slide, rx, ty, pw, ph, "b", "Adaptive Policy")
    add_panel(slide, lx, by, pw, ph, "c", "Image-to-Observation")
    add_panel(slide, rx, by, pw, ph, "d", "Real-World Rehabilitation Platform")

    # Inter-panel arrows.
    add_arrow(slide, lx + pw - 0.03, ty + 1.22, 0.58, 0.12, "right", C["blue"])
    add_text(slide, lx + pw - 0.08, ty + 1.03, 0.68, 0.18, "O_t, r_t", size=8.5, color=C["blue"])
    add_arrow(slide, lx + pw - 0.03, ty + 1.76, 0.58, 0.12, "left", C["orange"])
    add_text(slide, lx + pw + 0.19, ty + 1.91, 0.30, 0.18, "a_t", size=9, color=C["orange"])

    add_arrow(slide, rx + 2.05, ty + ph - 0.01, 0.12, 0.58, "down", C["green"])
    add_text(slide, rx + 2.12, ty + ph + 0.22, 0.92, 0.22, "trained policy", size=8.5, color=C["green"], align=PP_ALIGN.LEFT)

    add_arrow(slide, lx + 2.10, ty + ph - 0.01, 0.12, 0.58, "up", C["purple"])
    add_text(slide, lx + 1.43, ty + ph + 0.22, 0.72, 0.20, "same task", size=8.5, color=C["purple"])

    add_arrow(slide, lx + pw - 0.03, by + 1.16, 0.58, 0.12, "left", C["green"])
    add_text(slide, lx + pw + 0.03, by + 0.96, 0.52, 0.18, "image feed", size=8.5, color=C["green"])
    add_arrow(slide, lx + pw - 0.03, by + 1.85, 0.58, 0.12, "right", C["blue"])
    add_text(slide, lx + pw + 0.02, by + 2.02, 0.58, 0.18, "agent O_t", size=8.5, color=C["blue"])

    # Panel a: task world plus CMD-DR summary.
    world_x, world_y, world_w, world_h = lx + 0.38, ty + 0.92, 2.15, 1.62
    add_box(slide, world_x, world_y, world_w, world_h, fill=C["light_gray"], line=C["mid_gray"], line_width=1.0, radius=False)
    for i in range(1, 7):
        add_line(slide, world_x + i * world_w / 8, world_y, world_x + i * world_w / 8, world_y + world_h, color=rgb("#dde3ea"), width=0.25)
    for i in range(1, 5):
        add_line(slide, world_x, world_y + i * world_h / 6, world_x + world_w, world_y + i * world_h / 6, color=rgb("#dde3ea"), width=0.25)
    add_dot(slide, world_x + 1.02, world_y + 0.54, 0.48, C["green_light"], line=C["green"])
    add_text(slide, world_x + 1.00, world_y + 0.43, 0.50, 0.18, "ZPD", size=7.1, color=C["green"])
    add_line(slide, world_x + 0.26, world_y + 0.48, world_x + 0.66, world_y + 0.62, C["blue"], width=1.0)
    add_line(slide, world_x + 0.66, world_y + 0.62, world_x + 1.06, world_y + 0.83, C["blue"], width=1.0)
    add_line(slide, world_x + 1.08, world_y + 0.94, world_x + 1.44, world_y + 1.12, C["orange"], width=1.0)
    add_dot(slide, world_x + 0.23, world_y + 0.42, 0.11, C["blue"])
    add_dot(slide, world_x + 1.42, world_y + 1.04, 0.18, C["orange"])
    add_dot(slide, world_x + 1.06, world_y + 0.80, 0.08, C["red"])
    add_text(slide, world_x + 0.04, world_y + world_h + 0.08, world_w - 0.08, 0.18, "ZPD pursuit task workspace", size=7.6, color=C["gray"])

    cmd_x, cmd_y = lx + 2.72, ty + 0.86
    add_text(slide, cmd_x, cmd_y - 0.20, 1.38, 0.18, "CMD-DR virtual patient", size=8.8, color=C["orange"], bold=True)
    add_box(slide, cmd_x, cmd_y + 0.04, 1.36, 0.45, "Intent\nscripted / learned", fill=C["orange_light"], line=C["orange"], size=7.7)
    add_box(slide, cmd_x, cmd_y + 0.62, 1.36, 0.52, "Motor execution\ndelay / noise / smoothing", fill=RGBColor(255, 255, 255), line=C["orange"], size=7.3)
    add_box(slide, cmd_x, cmd_y + 1.32, 1.36, 0.38, "Patient diversity", fill=C["purple_light"], line=C["purple"], dash=True, size=7.8, bold=True)
    add_line(slide, cmd_x + 0.68, cmd_y + 0.49, cmd_x + 0.68, cmd_y + 0.62, C["orange"], width=1.0)
    add_line(slide, cmd_x + 0.68, cmd_y + 1.14, cmd_x + 0.68, cmd_y + 1.32, C["purple"], width=1.0)

    add_box(slide, lx + 0.62, ty + 2.76, 1.07, 0.32, "Observation O_t", fill=C["blue_light"], line=C["blue"], size=7.9, bold=True)
    add_box(slide, lx + 1.86, ty + 2.76, 1.07, 0.32, "Action a_t", fill=C["orange_light"], line=C["orange"], size=7.9, bold=True)
    add_box(slide, lx + 3.10, ty + 2.76, 1.07, 0.32, "Reward r_t", fill=C["purple_light"], line=C["purple"], size=7.9, bold=True)

    # Panel b: abstract policy network.
    inner_x, inner_y = rx + 0.34, ty + 0.62
    inner_w, inner_h = 3.88, 2.22
    add_box(slide, inner_x, inner_y, inner_w, inner_h, fill=RGBColor(255, 255, 255), line=C["black"], line_width=1.25)
    add_text(slide, inner_x + 0.18, inner_y + 0.14, 0.78, 0.26, "Input layer", size=10.1, bold=True)
    add_text(slide, inner_x + 1.48, inner_y + 0.14, 1.18, 0.26, "Hidden layers", size=10.1, bold=True)
    add_text(slide, inner_x + 3.02, inner_y + 0.14, 0.76, 0.26, "Output layer", size=10.1, bold=True)

    input_x = inner_x + 0.46
    hidden1_x = inner_x + 1.42
    hidden2_x = inner_x + 2.45
    output_x = inner_x + 3.55
    visible_y = [inner_y + 0.68, inner_y + 1.05, inner_y + 1.58, inner_y + 1.92]
    input_y = [inner_y + 0.72, inner_y + 1.05, inner_y + 1.61, inner_y + 1.94]
    output_y = inner_y + 1.30
    node_d = 0.21
    hidden_d = 0.26

    # Draw dense abstract connections first so nodes remain visually on top.
    for iy in input_y:
        for hy in visible_y:
            add_line(
                slide,
                input_x + node_d,
                iy + node_d / 2,
                hidden1_x,
                hy + hidden_d / 2,
                color=rgb("#9bb7df"),
                width=0.52,
            )
    for hy1 in visible_y:
        for hy2 in visible_y:
            add_line(
                slide,
                hidden1_x + hidden_d,
                hy1 + hidden_d / 2,
                hidden2_x,
                hy2 + hidden_d / 2,
                color=rgb("#9bb7df"),
                width=0.52,
            )
    for hy in visible_y:
        add_line(
            slide,
            hidden2_x + hidden_d,
            hy + hidden_d / 2,
            output_x,
            output_y + hidden_d / 2,
            color=rgb("#9bb7df"),
            width=0.52,
        )

    for iy in input_y:
        add_dot(slide, input_x, iy, node_d, C["orange_light"], line=C["orange"])
    for hy in visible_y:
        add_dot(slide, hidden1_x, hy, hidden_d, C["blue_light"], line=C["blue"])
        add_dot(slide, hidden2_x, hy, hidden_d, C["blue_light"], line=C["blue"])
    for x in [hidden1_x + 0.10, hidden2_x + 0.10]:
        for dy in [inner_y + 1.33, inner_y + 1.49]:
            add_dot(slide, x, dy, 0.060, C["black"])
    add_dot(slide, output_x, output_y, 0.28, C["green_light"], line=C["green"])

    add_text(slide, inner_x + 0.07, inner_y + 0.42, 0.82, 0.18, "spatial cues", size=8.0)
    add_text(slide, inner_x + 0.03, inner_y + 1.28, 0.90, 0.18, "temporal cues", size=8.0)
    add_text(slide, output_x - 0.18, output_y + 0.50, 0.82, 0.20, "action a_t", size=8.2)

    # Panel c: image processing into the observation consumed by the agent.
    cam_x, cam_y = lx + 0.58, by + 1.28
    add_box(slide, cam_x, cam_y, 1.20, 1.05, fill=C["light_gray"], line=C["mid_gray"], line_width=1.1, radius=False, dash=True)
    add_box(slide, cam_x + 0.23, cam_y + 0.63, 0.34, 0.22, fill=RGBColor(255, 255, 255), line=C["blue"], radius=False)
    add_dot(slide, cam_x + 0.30, cam_y + 0.67, 0.14, C["blue"])
    add_box(slide, cam_x + 0.68, cam_y + 0.30, 0.37, 0.25, fill=RGBColor(255, 255, 255), line=C["orange"], radius=False)
    add_dot(slide, cam_x + 0.77, cam_y + 0.34, 0.18, C["orange"])
    add_text(slide, cam_x + 0.08, cam_y + 1.12, 1.04, 0.18, "camera image", size=7.6, color=C["gray"])
    add_arrow(slide, lx + 1.92, by + 1.72, 0.30, 0.12, "right", C["green"])

    proc_x = lx + 2.28
    add_box(
        slide,
        proc_x,
        by + 1.18,
        1.05,
        1.24,
        "Detect + map\n\nmicrorobot p_m\nhand p_h\nworkspace coords",
        fill=C["green_light"],
        line=C["green"],
        size=7.2,
        bold=True,
    )
    add_arrow(slide, lx + 3.38, by + 1.72, 0.30, 0.12, "right", C["blue"])

    obs_x = lx + 3.74
    add_box(
        slide,
        obs_x,
        by + 1.18,
        0.72,
        1.24,
        "O_t\n\npositions\ndistance\nhistory",
        fill=C["blue_light"],
        line=C["blue"],
        size=7.0,
        bold=True,
    )

    # Panel d: hardware topology around the real platform placeholder.
    photo_x, photo_y = rx + 1.02, by + 1.17
    add_placeholder(slide, photo_x, photo_y, 1.82, 1.23, "insert real\nsetup photo")

    cam_node = add_box(slide, rx + 0.45, by + 1.48, 0.92, 0.38, "camera\n+ hand", fill=C["green_light"], line=C["green"], size=7.3)
    ctrl_node = add_box(slide, rx + 3.00, by + 0.88, 1.06, 0.40, "20 Hz policy\ninference", fill=C["blue_light"], line=C["blue"], size=7.4)
    ur_node = add_box(slide, rx + 3.06, by + 1.62, 0.94, 0.38, "UR10\nRTDE", fill=C["green_light"], line=C["green"], size=7.5)
    mag_node = add_box(slide, rx + 2.86, by + 2.50, 1.14, 0.40, "magnetic\nactuation", fill=C["orange_light"], line=C["orange"], size=7.3)
    ws_node = add_placeholder(slide, rx + 0.50, by + 2.50, 1.04, 0.40, "workspace\ninset")
    add_line(slide, rx + 1.37, by + 1.67, photo_x, photo_y + 0.56, color=C["green"], width=0.9)
    add_line(slide, rx + 2.84, by + 1.76, rx + 3.06, by + 1.80, color=C["green"], width=0.9)
    add_line(slide, rx + 2.63, by + 2.16, rx + 2.86, by + 2.58, color=C["orange"], width=0.9)
    add_line(slide, rx + 1.54, by + 2.70, photo_x + 0.38, photo_y + 1.10, color=C["mid_gray"], width=0.9, dash=True)
    add_arrow(slide, rx + 0.98, by + 1.12, 2.00, 0.10, "right", C["green"])
    add_arrow(slide, rx + 3.45, by + 1.31, 0.10, 0.28, "down", C["blue"])
    add_arrow(slide, rx + 3.43, by + 2.03, 0.10, 0.36, "down", C["orange"])

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
