from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[3]
OUT_DIR = ROOT / "manuscripts" / "figures" / "paper_ready"
OUT = OUT_DIR / "fig_dual_stream_gru_aux_architecture.pptx"


def I(value: float):
    return Inches(value)


COL = {
    "ink": RGBColor(34, 39, 46),
    "muted": RGBColor(102, 112, 133),
    "grid": RGBColor(223, 230, 238),
    "spatial": RGBColor(45, 95, 138),
    "spatial_fill": RGBColor(232, 242, 250),
    "spatial_fill_2": RGBColor(209, 229, 244),
    "temporal": RGBColor(132, 104, 48),
    "temporal_fill": RGBColor(255, 246, 219),
    "temporal_fill_2": RGBColor(244, 218, 158),
    "fusion": RGBColor(37, 128, 126),
    "fusion_fill": RGBColor(229, 246, 244),
    "ppo": RGBColor(63, 78, 116),
    "ppo_fill": RGBColor(236, 240, 252),
    "aux": RGBColor(155, 59, 76),
    "aux_fill": RGBColor(252, 235, 238),
    "green": RGBColor(57, 139, 73),
    "green_fill": RGBColor(230, 246, 232),
    "white": RGBColor(255, 255, 255),
}


def set_text(shape, text, size=10, color=None, bold=False, align=PP_ALIGN.CENTER):
    color = color or COL["ink"]
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    first = True
    for line in str(text).split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        run = p.add_run()
        run.text = line
        run.font.name = "Arial"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def add_text(slide, x, y, w, h, text, size=10, color=None, bold=False, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    set_text(shape, text, size=size, color=color, bold=bold, align=align)
    return shape


def add_box(
    slide,
    x,
    y,
    w,
    h,
    text="",
    fill=None,
    line=None,
    size=10,
    bold=False,
    radius=True,
    dashed=False,
    lw=1.5,
    text_color=None,
    align=PP_ALIGN.CENTER,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, I(x), I(y), I(w), I(h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or COL["grid"]
    shape.line.width = Pt(lw)
    if dashed:
        shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if text:
        set_text(shape, text, size=size, color=text_color, bold=bold, align=align)
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=None, width=2.0, dashed=False, begin=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
    conn.line.color.rgb = color or COL["ink"]
    conn.line.width = Pt(width)
    ln = conn.line._get_or_add_ln()
    tail = ln.find(qn("a:tailEnd"))
    if tail is None:
        tail = OxmlElement("a:tailEnd")
        ln.append(tail)
    tail.set("type", "triangle")
    if begin:
        head = ln.find(qn("a:headEnd"))
        if head is None:
            head = OxmlElement("a:headEnd")
            ln.append(head)
        head.set("type", "triangle")
    if dashed:
        conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return conn


def add_line(slide, x1, y1, x2, y2, color=None, width=1.0, dashed=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
    conn.line.color.rgb = color or COL["ink"]
    conn.line.width = Pt(width)
    if dashed:
        conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return conn


def add_dot(slide, x, y, r, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, I(x), I(y), I(r), I(r))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(1.0)
    return shape


def add_module_frame(slide, x, y, w, h, title, color):
    frame = add_box(slide, x, y, w, h, fill=RGBColor(253, 254, 255), line=color, radius=False, dashed=True, lw=1.2)
    add_text(slide, x + 0.12, y + 0.06, w - 0.24, 0.28, title, size=12, color=color, bold=True)
    return frame


def add_stack(slide, x, y, w, h, fill, line, text, size=8.5):
    for idx in range(3, 0, -1):
        dx = 0.07 * idx
        dy = 0.05 * idx
        add_box(slide, x + dx, y - dy, w, h, fill=RGBColor(248, 251, 255), line=line, radius=False, lw=0.8)
    top = add_box(slide, x, y, w, h, text, fill=fill, line=line, size=size, bold=True, radius=False, lw=1.2)
    return top


def add_vector_icon(slide, x, y, w, h, color):
    add_box(slide, x, y, w, h, fill=COL["white"], line=color, radius=False, lw=1.1)
    for idx, length in enumerate([0.62, 0.88, 0.46, 0.75, 0.55]):
        yy = y + 0.15 + idx * 0.12
        add_box(slide, x + 0.16, yy, w * length, 0.035, fill=color, line=color, radius=False, lw=0.4)
    add_text(slide, x + 0.08, y + h - 0.23, w - 0.16, 0.16, "12-D", size=7.2, color=color, bold=True)


def add_history_cards(slide, x, y, w, h, color):
    offsets = [(0.24, -0.18), (0.12, -0.09), (0.0, 0.0)]
    for idx, (dx, dy) in enumerate(offsets):
        card = add_box(slide, x + dx, y + dy, w, h, fill=COL["white"], line=color, radius=False, lw=0.9)
        card.shadow.inherit = False
        robot_x = x + dx + w * 0.30 + idx * 0.02
        robot_y = y + dy + h * 0.62 - idx * 0.05
        hand_x = x + dx + w * 0.68 - idx * 0.03
        hand_y = y + dy + h * 0.34 + idx * 0.03
        add_dot(slide, robot_x, robot_y, 0.06, COL["spatial"])
        add_dot(slide, hand_x, hand_y, 0.06, COL["aux"])
        add_arrow(slide, robot_x + 0.03, robot_y + 0.03, hand_x + 0.03, hand_y + 0.03, color=color, width=1.0)
    add_text(slide, x + 0.05, y + h + 0.05, w, 0.18, "t-15 ... t", size=7.5, color=color, bold=True)


def add_small_fc(slide, x, y, w, h, title, sub, color, fill):
    add_box(slide, x, y, w, h, fill=fill, line=color, radius=True, lw=1.4)
    add_text(slide, x + 0.06, y + 0.08, w - 0.12, h * 0.38, title, size=9, color=COL["ink"], bold=True)
    add_text(slide, x + 0.06, y + h * 0.47, w - 0.12, h * 0.36, sub, size=7.4, color=color)


def build_deck():
    prs = Presentation()
    prs.slide_width = I(13.33)
    prs.slide_height = I(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COL["white"]

    add_text(
        slide,
        0.35,
        0.08,
        12.65,
        0.35,
        "Dual-stream PPO policy network with auxiliary future dynamics",
        size=15,
        color=COL["ink"],
        bold=True,
    )

    add_box(slide, 0.20, 0.48, 12.92, 6.58, fill=None, line=RGBColor(165, 174, 186), radius=False, dashed=True, lw=1.1)

    # Left encoder branches.
    add_module_frame(slide, 0.45, 0.82, 4.55, 2.40, "Spatial stream", COL["spatial"])
    add_module_frame(slide, 0.45, 3.52, 4.55, 2.45, "Temporal interaction-history stream", COL["temporal"])

    add_text(slide, 0.62, 0.56, 1.55, 0.22, "Observation o_t", size=9.5, color=COL["muted"], bold=True, align=PP_ALIGN.LEFT)
    add_line(slide, 0.96, 0.78, 0.96, 5.92, color=RGBColor(190, 198, 209), width=1.0, dashed=True)
    add_text(slide, 0.65, 3.10, 0.72, 0.20, "split", size=7.5, color=COL["muted"])

    add_vector_icon(slide, 0.78, 1.36, 0.85, 0.84, COL["spatial"])
    add_text(slide, 0.60, 2.28, 1.28, 0.38, "scalar state s_t\npositions, distance, boundary,\nstride, previous action", size=7.5, color=COL["spatial"])
    add_small_fc(slide, 2.05, 1.38, 1.05, 0.86, "MLP", "FC -> LN -> ReLU\nscalar_dim -> 32", COL["spatial"], COL["spatial_fill"])
    add_stack(slide, 3.60, 1.45, 0.92, 0.72, COL["spatial_fill_2"], COL["spatial"], "z_s\n32", size=9)
    add_arrow(slide, 1.70, 1.78, 2.02, 1.78, color=COL["spatial"], width=2.0)
    add_arrow(slide, 3.12, 1.78, 3.56, 1.78, color=COL["spatial"], width=2.0)

    add_history_cards(slide, 0.70, 4.18, 0.95, 0.68, COL["temporal"])
    add_text(slide, 0.55, 5.20, 1.35, 0.38, "16-frame history\nrelative interaction\nmotion sequence", size=7.3, color=COL["temporal"])
    add_small_fc(slide, 2.05, 4.14, 1.05, 0.92, "GRU", "C -> 64\nlast hidden state", COL["temporal"], COL["temporal_fill"])
    add_small_fc(slide, 3.42, 4.14, 1.12, 0.92, "Projection", "LN + FC + ReLU\n64 -> 32", COL["temporal"], COL["temporal_fill_2"])
    add_arrow(slide, 1.72, 4.55, 2.02, 4.55, color=COL["temporal"], width=2.0)
    add_arrow(slide, 3.12, 4.55, 3.39, 4.55, color=COL["temporal"], width=2.0)

    # Fusion module.
    add_module_frame(slide, 5.22, 2.38, 2.35, 2.05, "Fusion encoder", COL["fusion"])
    add_box(slide, 5.40, 3.04, 0.76, 0.58, "Concat\nz_s + z_h", fill=COL["fusion_fill"], line=COL["fusion"], size=7.2, bold=True, lw=1.3)
    add_small_fc(slide, 6.28, 2.94, 0.86, 0.78, "Fusion MLP", "LN + FC\n64 -> 64", COL["fusion"], COL["fusion_fill"])
    add_box(slide, 6.24, 3.86, 0.90, 0.34, "shared feature f_t", fill=RGBColor(214, 239, 237), line=COL["fusion"], size=7.5, bold=True, lw=1.2)
    add_arrow(slide, 4.54, 1.78, 5.40, 3.15, color=COL["spatial"], width=2.0)
    add_arrow(slide, 4.56, 4.55, 5.40, 3.50, color=COL["temporal"], width=2.0)
    add_arrow(slide, 6.18, 3.33, 6.25, 3.33, color=COL["fusion"], width=1.6)
    add_arrow(slide, 6.68, 3.73, 6.68, 3.84, color=COL["fusion"], width=1.4)

    # Auxiliary module.
    add_module_frame(slide, 7.90, 0.82, 4.85, 2.48, "Auxiliary future-dynamics module", COL["aux"])
    add_small_fc(slide, 8.15, 1.62, 0.86, 0.72, "Shared FC", "64 -> 64", COL["aux"], COL["aux_fill"])
    add_small_fc(slide, 9.26, 1.18, 1.08, 0.70, "Future traj.", "64 -> H*2", COL["aux"], RGBColor(248, 222, 226))
    add_small_fc(slide, 9.26, 2.20, 1.08, 0.70, "Risk head", "64 -> 32 -> 1", COL["aux"], RGBColor(248, 222, 226))
    add_box(slide, 10.70, 1.12, 0.98, 0.82, "predicted\nfuture hand\nmotion\nH=8", fill=COL["white"], line=COL["aux"], size=7.6, lw=1.2)
    add_box(slide, 10.70, 2.18, 0.98, 0.70, "near-catch\nrisk logit", fill=COL["white"], line=COL["aux"], size=7.6, lw=1.2)
    add_box(slide, 11.86, 1.56, 0.70, 0.86, "rollout\ntargets\nD, y", fill=RGBColor(255, 247, 248), line=COL["aux"], size=7.2, dashed=True, lw=1.0)
    add_arrow(slide, 7.14, 3.98, 8.13, 1.98, color=COL["aux"], width=1.8)
    add_arrow(slide, 9.02, 1.98, 9.24, 1.52, color=COL["aux"], width=1.5)
    add_arrow(slide, 9.02, 1.98, 9.24, 2.55, color=COL["aux"], width=1.5)
    add_arrow(slide, 10.36, 1.53, 10.67, 1.53, color=COL["aux"], width=1.5)
    add_arrow(slide, 10.36, 2.55, 10.67, 2.55, color=COL["aux"], width=1.5)
    add_arrow(slide, 11.84, 1.86, 11.67, 1.58, color=COL["aux"], width=1.0, dashed=True, begin=True)
    add_text(slide, 8.20, 2.92, 3.92, 0.22, "aux loss: L_traj + L_risk updates the shared encoder", size=7.8, color=COL["aux"])

    # PPO module.
    add_module_frame(slide, 7.90, 3.72, 4.85, 2.25, "PPO actor-critic heads", COL["ppo"])
    add_small_fc(slide, 8.34, 4.35, 0.90, 0.62, "Actor", "policy MLP", COL["ppo"], COL["ppo_fill"])
    add_small_fc(slide, 8.34, 5.18, 0.90, 0.62, "Critic", "value MLP", RGBColor(91, 76, 119), RGBColor(244, 239, 251))
    add_box(slide, 9.70, 4.35, 1.08, 0.62, "2D action\na_t=(dx,dy)", fill=COL["white"], line=COL["ppo"], size=8.0, lw=1.2)
    add_box(slide, 9.70, 5.18, 1.08, 0.62, "state value\nV(o_t)", fill=COL["white"], line=RGBColor(91, 76, 119), size=8.0, lw=1.2)
    add_box(slide, 11.15, 4.64, 1.05, 0.72, "rollout\nbuffer\n(o,a,r,V)", fill=RGBColor(247, 249, 255), line=COL["ppo"], size=7.3, dashed=True, lw=1.0)
    add_box(slide, 11.15, 5.52, 1.05, 0.32, "PPO loss", fill=COL["ppo_fill"], line=COL["ppo"], size=7.8, bold=True, lw=1.0)
    add_arrow(slide, 7.14, 4.03, 8.30, 4.66, color=COL["ppo"], width=1.8)
    add_arrow(slide, 7.14, 4.03, 8.30, 5.48, color=RGBColor(91, 76, 119), width=1.8)
    add_arrow(slide, 9.26, 4.66, 9.67, 4.66, color=COL["ppo"], width=1.6)
    add_arrow(slide, 9.26, 5.49, 9.67, 5.49, color=RGBColor(91, 76, 119), width=1.6)
    add_arrow(slide, 11.68, 5.38, 11.68, 5.50, color=COL["ppo"], width=1.0)
    add_arrow(slide, 11.15, 5.68, 9.28, 5.68, color=COL["ppo"], width=1.0, dashed=True)

    # Shared training objective.
    add_box(
        slide,
        2.00,
        6.36,
        9.40,
        0.42,
        "Training objective: L_total = L_PPO + lambda_traj L_traj + lambda_risk L_risk   |   shared encoder receives PPO and auxiliary gradients",
        fill=RGBColor(249, 250, 252),
        line=RGBColor(185, 194, 205),
        size=8.1,
        radius=True,
        lw=1.0,
        text_color=COL["muted"],
    )
    add_arrow(slide, 6.68, 6.34, 6.68, 4.22, color=COL["fusion"], width=1.0, dashed=True)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build_deck()
