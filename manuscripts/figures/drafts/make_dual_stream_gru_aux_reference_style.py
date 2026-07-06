from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[3]
OUT_DIR = ROOT / "manuscripts" / "figures" / "paper_ready"
OUT = OUT_DIR / "fig_dual_stream_gru_aux_architecture_reference_style.pptx"


def I(v: float):
    return Inches(v)


C = {
    "ink": RGBColor(28, 32, 38),
    "muted": RGBColor(100, 110, 126),
    "grid": RGBColor(178, 188, 202),
    "blue": RGBColor(46, 108, 164),
    "blue_light": RGBColor(220, 238, 250),
    "blue_mid": RGBColor(185, 216, 240),
    "orange": RGBColor(211, 120, 50),
    "orange_light": RGBColor(255, 241, 224),
    "orange_mid": RGBColor(248, 209, 166),
    "teal": RGBColor(39, 132, 130),
    "teal_light": RGBColor(224, 246, 244),
    "green": RGBColor(111, 178, 83),
    "green_light": RGBColor(225, 245, 218),
    "yellow": RGBColor(230, 188, 62),
    "yellow_light": RGBColor(255, 246, 203),
    "red": RGBColor(209, 42, 48),
    "red_light": RGBColor(255, 231, 233),
    "navy": RGBColor(68, 82, 123),
    "navy_light": RGBColor(232, 238, 250),
    "purple": RGBColor(92, 78, 126),
    "purple_light": RGBColor(241, 235, 250),
    "white": RGBColor(255, 255, 255),
}


def add_text(slide, x, y, w, h, s, size=7.0, color=None, bold=False, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    first = True
    for line in str(s).split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        r = p.add_run()
        r.text = line
        r.font.name = "Arial"
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color or C["ink"]
    return box


def rect(slide, x, y, w, h, fill=None, line=None, lw=0.8, dashed=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(x), I(y), I(w), I(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line or C["ink"]
    shp.line.width = Pt(lw)
    if dashed:
        shp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return shp


def ellipse(slide, x, y, w, h, fill=None, line_color=None, lw=0.7, text="", color=None, size=6.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, I(x), I(y), I(w), I(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill or C["white"]
    shp.line.color.rgb = line_color or C["ink"]
    shp.line.width = Pt(lw)
    if text:
        add_text(slide, x, y, w, h, text, size=size, color=color or line_color or C["ink"], bold=True)
    return shp


def connector(slide, x1, y1, x2, y2, color=None, width=0.75, dashed=False, arrow=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
    conn.line.color.rgb = color or C["ink"]
    conn.line.width = Pt(width)
    if dashed:
        conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if arrow:
        ln = conn.line._get_or_add_ln()
        tail = ln.find(qn("a:tailEnd"))
        if tail is None:
            tail = OxmlElement("a:tailEnd")
            ln.append(tail)
        tail.set("type", "triangle")
    return conn


def ortho(slide, pts, color=None, width=0.75, dashed=False, arrow=True):
    for i in range(len(pts) - 1):
        connector(
            slide,
            pts[i][0],
            pts[i][1],
            pts[i + 1][0],
            pts[i + 1][1],
            color=color,
            width=width,
            dashed=dashed,
            arrow=arrow and i == len(pts) - 2,
        )


def tensor(slide, x, y, w, h, label, dim, color, fill, n=4, dx=0.045, dy=-0.035, slices=True):
    for i in range(n - 1, 0, -1):
        rect(slide, x + dx * i, y + dy * i, w, h, fill=C["white"], line=color, lw=0.45)
    rect(slide, x, y, w, h, fill=fill, line=color, lw=0.9)
    if slices:
        for j in range(1, 4):
            yy = y + h * j / 4
            connector(slide, x + 0.03, yy, x + w - 0.03, yy, color=RGBColor(206, 216, 228), width=0.35)
    add_text(slide, x - 0.16, y + h + 0.04, w + 0.32, 0.16, label, size=6.6, color=color, bold=True)
    if dim:
        add_text(slide, x - 0.14, y + h + 0.22, w + 0.28, 0.13, dim, size=5.9, color=C["muted"])


def vector(slide, x, y, w, h, label, dim, color, fill, slots=6):
    rect(slide, x, y, w, h, fill=fill, line=color, lw=0.9)
    for j in range(1, slots):
        yy = y + h * j / slots
        connector(slide, x, yy, x + w, yy, color=RGBColor(210, 220, 232), width=0.35)
    add_text(slide, x - 0.18, y + h + 0.04, w + 0.36, 0.15, label, size=6.5, color=color, bold=True)
    add_text(slide, x - 0.18, y + h + 0.20, w + 0.36, 0.12, dim, size=5.7, color=C["muted"])


def bar(slide, x, y, w, h, label, dim, color, fill):
    rect(slide, x, y, w, h, fill=fill, line=color, lw=0.85)
    if label:
        add_text(slide, x - 0.18, y - 0.19, w + 0.36, 0.15, label, size=6.2, color=color, bold=True)
    if dim:
        add_text(slide, x - 0.24, y + h + 0.04, w + 0.48, 0.12, dim, size=5.6, color=C["muted"])


def gru_unroll(slide, x, y, color, fill):
    xs = [x, x + 0.57, x + 1.14]
    labels = ["x_{t-15}", "...", "x_t"]
    for i, xx in enumerate(xs):
        rect(slide, xx + 0.08, y + 0.60, 0.12, 0.38, fill=fill, line=color, lw=0.7)
        for j in range(1, 3):
            yy = y + 0.60 + 0.38 * j / 3
            connector(slide, xx + 0.08, yy, xx + 0.20, yy, color=RGBColor(226, 205, 180), width=0.3)
        add_text(slide, xx - 0.07, y + 0.99, 0.42, 0.12, labels[i], size=5.0, color=color)
        rect(slide, xx, y, 0.28, 0.42, fill=fill, line=color, lw=0.85)
        add_text(slide, xx + 0.01, y + 0.13, 0.26, 0.12, "GRU", size=5.0, color=color, bold=True)
        connector(slide, xx + 0.14, y + 0.60, xx + 0.14, y + 0.44, color=color, width=0.55, arrow=True)
        if i < 2:
            connector(slide, xx + 0.28, y + 0.21, xs[i + 1] - 0.04, y + 0.21, color=color, width=0.75, arrow=True)
            add_text(slide, xx + 0.31, y + 0.04, 0.26, 0.10, "h", size=4.5, color=color)
    connector(slide, xs[2] + 0.28, y + 0.21, xs[2] + 0.42, y + 0.21, color=color, width=0.75, arrow=True)
    add_text(slide, x + 0.04, y - 0.25, 1.34, 0.15, "unrolled GRU over time", size=5.8, color=color, bold=True)


def build():
    prs = Presentation()
    prs.slide_width = I(13.33)
    prs.slide_height = I(5.80)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C["white"]

    # Outer frame and title.
    rect(slide, 0.16, 0.18, 13.00, 5.42, fill=None, line=RGBColor(105, 109, 116), lw=0.9, dashed=True)

    # Module frames, close to the reference figure.
    rect(slide, 0.42, 0.52, 4.85, 2.05, fill=None, line=C["blue"], lw=0.75, dashed=True)
    rect(slide, 0.42, 3.00, 4.85, 2.00, fill=None, line=C["orange"], lw=0.75, dashed=True)
    rect(slide, 5.52, 0.52, 2.35, 4.48, fill=None, line=C["ink"], lw=0.75, dashed=True)
    rect(slide, 7.98, 0.52, 4.85, 2.06, fill=None, line=C["red"], lw=0.75, dashed=True)
    rect(slide, 7.98, 2.93, 4.85, 2.07, fill=None, line=RGBColor(150, 156, 168), lw=0.75, dashed=True)
    add_text(slide, 1.78, 0.58, 2.10, 0.16, "Spatial state branch", 7.2, C["blue"], True)
    add_text(slide, 1.34, 3.06, 2.90, 0.16, "Temporal interaction-history branch", 7.2, C["orange"], True)
    add_text(slide, 5.90, 0.58, 1.60, 0.16, "Fusion module", 7.0, C["ink"], True)
    add_text(slide, 9.25, 0.58, 2.40, 0.16, "Auxiliary task module", 7.0, C["red"], True)
    add_text(slide, 9.20, 2.99, 2.45, 0.16, "PPO actor-critic module", 7.0, C["navy"], True)

    # Observation split.
    add_text(slide, 0.55, 0.30, 1.12, 0.14, "Observation o_t", 5.8, C["muted"], True, PP_ALIGN.LEFT)
    connector(slide, 0.96, 0.44, 0.96, 1.02, color=C["grid"], width=0.5, dashed=True, arrow=True)
    connector(slide, 0.96, 0.44, 0.96, 3.50, color=C["grid"], width=0.5, dashed=True, arrow=True)

    # Top branch: scalar stream.
    vector(slide, 0.78, 1.10, 0.22, 0.96, "s_t", "12 x 1", C["blue"], C["blue_light"], slots=6)
    bar(slide, 1.64, 1.12, 0.15, 0.92, "FC", "12 -> 32", C["blue"], C["blue_mid"])
    bar(slide, 2.18, 1.12, 0.15, 0.92, "Norm", "", C["blue"], C["blue_light"])
    tensor(slide, 2.88, 1.16, 0.36, 0.84, "spatial emb.", "32 x 1", C["blue"], C["blue_mid"], n=4)
    bar(slide, 3.74, 1.22, 0.12, 0.72, "Flat", "", C["blue"], C["blue_light"])
    ortho(slide, [(1.00, 1.58), (1.62, 1.58)], C["blue"], 0.9, arrow=True)
    ortho(slide, [(1.80, 1.58), (2.16, 1.58)], C["blue"], 0.8, arrow=True)
    ortho(slide, [(2.34, 1.58), (2.86, 1.58)], C["blue"], 0.8, arrow=True)
    ortho(slide, [(3.25, 1.58), (3.72, 1.58)], C["blue"], 0.8, arrow=True)

    # Bottom branch: temporal stream.
    tensor(slide, 0.70, 3.57, 0.42, 0.82, "h_{t-15:t}", "16 x C", C["orange"], C["orange_light"], n=5)
    gru_unroll(slide, 1.58, 3.66, C["orange"], C["orange_light"])
    bar(slide, 3.08, 3.52, 0.22, 0.92, "h_T", "64", C["orange"], C["orange_mid"])
    bar(slide, 3.76, 3.57, 0.15, 0.82, "FC", "64 -> 32", C["orange"], C["orange_mid"])
    tensor(slide, 4.32, 3.65, 0.28, 0.64, "z_h", "32", C["orange"], C["orange_mid"], n=3)
    ortho(slide, [(1.13, 3.98), (1.56, 3.98)], C["orange"], 0.85, arrow=True)
    ortho(slide, [(3.04, 3.87), (3.06, 3.98)], C["orange"], 0.8, arrow=True)
    ortho(slide, [(3.31, 3.98), (3.74, 3.98)], C["orange"], 0.8, arrow=True)
    ortho(slide, [(3.92, 3.98), (4.30, 3.98)], C["orange"], 0.8, arrow=True)

    # Fusion module with rectilinear inputs.
    bar(slide, 5.74, 1.35, 0.12, 0.74, "FC", "32", C["blue"], C["blue_light"])
    bar(slide, 5.74, 3.55, 0.12, 0.74, "FC", "32", C["orange"], C["orange_light"])
    ellipse(slide, 6.08, 2.42, 0.20, 0.20, fill=C["white"], line_color=C["ink"], text="+", size=6.2)
    add_text(slide, 5.78, 2.66, 0.80, 0.12, "concat", 5.5, C["muted"])
    tensor(slide, 6.54, 1.92, 0.18, 1.20, "z", "64", C["green"], C["green_light"], n=1, slices=True)
    bar(slide, 7.12, 1.72, 0.22, 1.60, "FC", "64 -> 64", C["green"], C["green_light"])
    tensor(slide, 7.58, 1.86, 0.16, 1.32, "f_t", "64", C["yellow"], C["yellow_light"], n=1, slices=True)

    ortho(slide, [(3.86, 1.58), (5.48, 1.58), (5.48, 1.72), (5.72, 1.72)], C["blue"], 0.75, arrow=True)
    ortho(slide, [(4.60, 3.98), (5.48, 3.98), (5.48, 3.92), (5.72, 3.92)], C["orange"], 0.75, arrow=True)
    ortho(slide, [(5.86, 1.72), (6.02, 1.72), (6.02, 2.52), (6.08, 2.52)], C["blue"], 0.65, arrow=True)
    ortho(slide, [(5.86, 3.92), (6.02, 3.92), (6.02, 2.52), (6.08, 2.52)], C["orange"], 0.65, arrow=True)
    ortho(slide, [(6.28, 2.52), (6.52, 2.52)], C["green"], 0.85, arrow=True)
    ortho(slide, [(6.72, 2.52), (7.10, 2.52)], C["green"], 0.85, arrow=True)
    ortho(slide, [(7.34, 2.52), (7.56, 2.52)], C["yellow"], 0.85, arrow=True)

    # Shared feature bus, all orthogonal.
    connector(slide, 7.78, 2.52, 8.10, 2.52, color=C["ink"], width=0.65)
    connector(slide, 8.10, 1.42, 8.10, 4.30, color=C["ink"], width=0.65)

    # Auxiliary module.
    bar(slide, 8.42, 1.18, 0.16, 0.92, "FC", "64 -> 64", C["red"], C["red_light"])
    bar(slide, 9.18, 0.98, 0.16, 0.82, "FC", "64 -> H*2", C["red"], C["red_light"])
    tensor(slide, 9.70, 0.98, 0.40, 0.74, "D_hat", "H x 2", C["red"], C["red_light"], n=5)
    bar(slide, 9.18, 1.95, 0.16, 0.54, "FC", "64 -> 1", C["red"], C["red_light"])
    ellipse(slide, 9.88, 2.13, 0.13, 0.13, fill=C["red_light"], line_color=C["red"])
    add_text(slide, 9.75, 2.28, 0.42, 0.12, "risk", 5.4, C["red"])
    tensor(slide, 11.10, 1.02, 0.28, 0.70, "targets", "D, y", C["red"], C["white"], n=3)
    add_text(slide, 8.78, 2.39, 1.75, 0.12, "aux loss", 5.5, C["red"])
    ortho(slide, [(8.10, 1.62), (8.40, 1.62)], C["red"], 0.75, arrow=True)
    ortho(slide, [(8.58, 1.62), (8.92, 1.62), (8.92, 1.39), (9.16, 1.39)], C["red"], 0.75, arrow=True)
    ortho(slide, [(8.92, 1.62), (8.92, 2.22), (9.16, 2.22)], C["red"], 0.75, arrow=True)
    ortho(slide, [(9.34, 1.39), (9.68, 1.39)], C["red"], 0.75, arrow=True)
    ortho(slide, [(9.34, 2.22), (9.86, 2.22)], C["red"], 0.75, arrow=True)
    ortho(slide, [(11.08, 1.37), (10.24, 1.37)], C["red"], 0.55, dashed=True, arrow=True)

    # PPO module.
    bar(slide, 8.42, 3.55, 0.16, 0.76, "Actor", "", C["navy"], C["navy_light"])
    vector(slide, 9.18, 3.68, 0.14, 0.48, "a_t", "2", C["navy"], C["white"], slots=2)
    bar(slide, 8.42, 4.38, 0.16, 0.48, "Critic", "", C["purple"], C["purple_light"])
    ellipse(slide, 9.19, 4.52, 0.12, 0.12, fill=C["purple_light"], line_color=C["purple"])
    add_text(slide, 9.02, 4.68, 0.45, 0.12, "V_t", 5.4, C["purple"])
    tensor(slide, 10.92, 3.73, 0.32, 0.62, "rollout", "o,a,r,V", C["navy"], C["white"], n=3)
    add_text(slide, 10.50, 4.68, 1.10, 0.13, "PPO update", 5.8, C["navy"], True)
    ortho(slide, [(8.10, 3.94), (8.40, 3.94)], C["navy"], 0.75, arrow=True)
    ortho(slide, [(8.10, 4.62), (8.40, 4.62)], C["purple"], 0.75, arrow=True)
    ortho(slide, [(8.58, 3.94), (9.16, 3.94)], C["navy"], 0.75, arrow=True)
    ortho(slide, [(8.58, 4.62), (9.18, 4.58)], C["purple"], 0.75, arrow=True)
    ortho(slide, [(9.34, 3.94), (10.90, 3.94)], C["navy"], 0.55, dashed=True, arrow=True)
    ortho(slide, [(11.10, 4.36), (11.10, 4.58)], C["navy"], 0.55, dashed=True, arrow=True)
    ortho(slide, [(10.48, 4.74), (8.64, 4.74)], C["navy"], 0.55, dashed=True, arrow=True)

    # Training objective: small, caption-like.
    add_text(
        slide,
        5.15,
        5.16,
        3.80,
        0.16,
        "L_total = L_PPO + lambda_traj L_traj + lambda_risk L_risk",
        5.6,
        C["muted"],
    )

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
